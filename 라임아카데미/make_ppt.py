# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── 컬러 팔레트 ──────────────────────────────
BG     = RGBColor(0x0D, 0x0D, 0x0D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xD4, 0xFF, 0x00)   # 라임
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0x99, 0x99, 0x99)
CARD   = RGBColor(0x1A, 0x1A, 0x1A)
CARD2  = RGBColor(0x22, 0x22, 0x22)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(sl, c=BG):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = c

def rect(sl, l, t, w, h, c, line=False):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c
    if line: s.line.color.rgb = c
    else:    s.line.fill.background()
    return s

def txt(sl, s, l, t, w, h, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b = sl.shapes.add_textbox(l, t, w, h)
    b.word_wrap = wrap
    tf = b.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = s
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    r.font.name = "Apple SD Gothic Neo"
    return b

def line(sl, l, t, w, h, c=GRAY):
    rect(sl, l, t, w, h, c)

def tag(sl, s, l, t, w=Inches(1.8)):
    rect(sl, l, t, w, Inches(0.32), ACCENT)
    txt(sl, s, l, t, w, Inches(0.32), size=11, bold=True,
        color=BG, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SL 01 · 타이틀
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)

# 좌측 세로 라인
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

# 회차 태그
txt(sl, "WEEK 01", Inches(0.75), Inches(0.55), Inches(3), Inches(0.4),
    size=12, bold=True, color=ACCENT)

# 메인 타이틀
txt(sl, "돈복사", Inches(0.75), Inches(1.2), Inches(9), Inches(2.2),
    size=110, bold=True, color=WHITE)

# 부제
txt(sl, "우리가 왜 돈을 못 버나", Inches(0.75), Inches(3.3), Inches(8), Inches(0.7),
    size=28, color=LGRAY)

# 하단 가로선
line(sl, Inches(0.75), Inches(4.3), Inches(11.5), Inches(0.04))

# 통제 아이디어
txt(sl, "통제 아이디어", Inches(0.75), Inches(4.55), Inches(5), Inches(0.4),
    size=11, color=ACCENT, bold=True)
txt(sl,
    "주식 수익은, 감(感)이 아닌\n검증된 조건을 반복 실행할 수 있을 때 일어난다.",
    Inches(0.75), Inches(5.0), Inches(9), Inches(1.2),
    size=20, color=LGRAY)

# 우측 하단 숫자 워터마크
txt(sl, "01", Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.8),
    size=130, bold=True, color=RGBColor(0x1A,0x1A,0x1A), align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════
# SL 02 · 6주 커리큘럼
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "6주 커리큘럼", Inches(0.75), Inches(0.55), Inches(8), Inches(0.65),
    size=36, bold=True, color=WHITE)
line(sl, Inches(0.75), Inches(1.3), Inches(11.5), Inches(0.03))

weeks = [
    ("01", "자각",  "우리가 왜 돈을 못 버나",   True),
    ("02", "탐색",  "좋은 종목의 조건이란?",      False),
    ("03", "한계",  "조건 찾기의 현실적 한계",    False),
    ("04", "복잡",  "타이밍까지 더해지면?",       False),
    ("05", "해방",  "도구가 있다면 달라질까?",    False),
    ("06", "출발",  "나의 반복 실행 시작",        False),
]
for i,(num,kw,title,on) in enumerate(weeks):
    t = Inches(1.5) + Inches(0.97)*i
    nc = ACCENT if on else GRAY
    tc = WHITE  if on else GRAY
    txt(sl, num,   Inches(0.75), t, Inches(0.7), Inches(0.75), size=22, bold=True, color=nc)
    txt(sl, f"[ {kw} ]", Inches(1.55), t+Inches(0.14), Inches(1.5), Inches(0.45),
        size=14, bold=on, color=nc)
    txt(sl, title, Inches(3.2), t+Inches(0.14), Inches(7), Inches(0.45),
        size=19, bold=on, color=tc)
    if on:
        rect(sl, Inches(10.5), t+Inches(0.12), Inches(1.8), Inches(0.46), ACCENT)
        txt(sl, "◀ 오늘", Inches(10.5), t+Inches(0.12), Inches(1.8), Inches(0.46),
            size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    if i < 5:
        line(sl, Inches(0.75), t+Inches(0.93), Inches(11.5), Inches(0.02),
             RGBColor(0x22,0x22,0x22))

# ════════════════════════════════════════════════════════
# SL 03 · 오프닝 — 솔직한 소개
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "OPENING", Inches(0.75), Inches(0.55), Inches(4), Inches(0.4),
    size=12, bold=True, color=ACCENT)
txt(sl, "먼저 드릴 말씀이 있어요.", Inches(0.75), Inches(1.1), Inches(10), Inches(0.75),
    size=38, bold=True, color=WHITE)
line(sl, Inches(0.75), Inches(1.95), Inches(11.5), Inches(0.03))

# 큰 따옴표
txt(sl, "“", Inches(0.55), Inches(2.0), Inches(1.2), Inches(1.2),
    size=90, bold=True, color=ACCENT)

lines_o = [
    ("저도 한때 감으로 샀다가 크게 잃은 적 있어요.", True),
    ("그 뒤로 조건을 찾기 시작했고,", False),
    ("그 과정이 너무 힘들어서 도구를 직접 만들었습니다.", False),
    ("", False),
    ("오늘 제 도구를 팔려는 게 아닙니다.", False),
    ("5주차에 보여드릴 건데 — 그때 '아, 이게 왜 필요한지' 느끼실 거예요.", False),
    ("", False),
    ("오늘은 우리가 왜 돈을 못 벌었는지부터 솔직하게 얘기해봅시다.", True),
]
txb = sl.shapes.add_textbox(Inches(1.4), Inches(2.3), Inches(11.0), Inches(4.0))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
for idx,(ln,em) in enumerate(lines_o):
    p = tf.paragraphs[0] if idx==0 else tf.add_paragraph()
    r = p.add_run(); r.text = ln
    r.font.size  = Pt(21 if em else 18)
    r.font.bold  = em
    r.font.color.rgb = WHITE if em else LGRAY
    r.font.name  = "Apple SD Gothic Neo"

txt(sl, "오늘 여러분이 이 모임을 만들어가는 거예요.",
    Inches(0.75), Inches(6.75), Inches(10), Inches(0.45),
    size=13, color=ACCENT)

# ════════════════════════════════════════════════════════
# SL 04 · 자가진단 카드
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "자가진단", Inches(0.75), Inches(0.55), Inches(10), Inches(0.45),
    size=13, bold=True, color=ACCENT)
txt(sl, "당신의 투자 습관은?", Inches(0.75), Inches(1.05), Inches(10), Inches(0.7),
    size=38, bold=True, color=WHITE)
txt(sl, "각 항목을 읽고 솔직하게  ○  △  ×  로 표시해주세요. 정답은 없습니다.",
    Inches(0.75), Inches(1.8), Inches(11), Inches(0.4),
    size=15, color=LGRAY)
line(sl, Inches(0.75), Inches(2.3), Inches(11.5), Inches(0.03))

qs = [
    "나는 지인, 커뮤니티 의견에 흔들리지 않고 매매한다.",
    "나는 투자를 하기 전 공부를 하고 투자를 한다.",
    "나는 내가 투자한 것에 근거 3가지 이상 댈 수 있다.",
    "수익이 나도 목표한 금액이 있고, 그 금액까지 기다린다.",
    "나는 나의 전체 자산이 매년 얼마나 불어나는지 알고 있다.",
]
for i,q in enumerate(qs):
    t = Inches(2.5) + Inches(0.98)*i
    # 번호
    rect(sl, Inches(0.75), t, Inches(0.55), Inches(0.55),
         ACCENT if i==0 else CARD2)
    txt(sl, str(i+1), Inches(0.75), t, Inches(0.55), Inches(0.55),
        size=18, bold=True,
        color=BG if i==0 else LGRAY, align=PP_ALIGN.CENTER)
    # 질문
    txt(sl, q, Inches(1.45), t+Inches(0.07), Inches(9.5), Inches(0.45),
        size=18, bold=(i==0),
        color=WHITE if i==0 else LGRAY)
    # 응답칸
    for j,(sym,sc) in enumerate([("○",RGBColor(0xD4,0xFF,0x00)),
                                   ("△",RGBColor(0xFF,0xFF,0xFF)),
                                   ("×",RGBColor(0x66,0x66,0x66))]):
        bl = Inches(11.3) + Inches(0.52)*j
        rect(sl, bl, t, Inches(0.48), Inches(0.55), CARD2)
        txt(sl, sym, bl, t, Inches(0.48), Inches(0.55),
            size=18, bold=True, color=sc, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SL 05 · 브리지 — 불편한 질문
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "잠깐 생각해봐요.", Inches(0.75), Inches(0.55), Inches(10), Inches(0.6),
    size=38, bold=True, color=WHITE)
line(sl, Inches(0.75), Inches(1.3), Inches(11.5), Inches(0.03))

qs2 = [
    ("Q1", "가장 크게 잃었던 매매를 떠올려보세요."),
    ("Q2", "그 종목을  '왜'  샀었나요?"),
    ("Q3", "그 이유를 지금 3가지 이상 댈 수 있나요?"),
]
for i,(label,q) in enumerate(qs2):
    t = Inches(1.9) + Inches(1.55)*i
    txt(sl, label, Inches(0.75), t, Inches(1.2), Inches(0.45),
        size=13, bold=True, color=ACCENT)
    txt(sl, q, Inches(0.75), t+Inches(0.45), Inches(11), Inches(0.8),
        size=30, bold=(i==1), color=WHITE if i==1 else LGRAY)

txt(sl, "이 질문들을 기억하면서 다음 활동을 시작합니다.",
    Inches(0.75), Inches(6.8), Inches(10), Inches(0.4),
    size=13, color=GRAY)

# ════════════════════════════════════════════════════════
# SL 06 · 활동 안내
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "본활동", Inches(0.75), Inches(0.55), Inches(5), Inches(0.45),
    size=13, bold=True, color=ACCENT)
txt(sl, "지금부터 할 것", Inches(0.75), Inches(1.05), Inches(9), Inches(0.65),
    size=38, bold=True, color=WHITE)
line(sl, Inches(0.75), Inches(1.8), Inches(11.5), Inches(0.03))

rules = [
    "발언 순서는 피스메이커가 정합니다.",
    "남의 매매를 평가하거나 비웃지 않습니다.",
    "틀린 답 없습니다. 솔직한 게 전부입니다.",
    "카드에 쓴 내용을 그대로 읽어도 됩니다.",
]
for i,r_ in enumerate(rules):
    t = Inches(2.1) + Inches(1.0)*i
    rect(sl, Inches(0.75), t+Inches(0.05), Inches(0.42), Inches(0.42), CARD2)
    txt(sl, str(i+1), Inches(0.75), t+Inches(0.05), Inches(0.42), Inches(0.42),
        size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(sl, r_, Inches(1.3), t, Inches(10), Inches(0.55),
        size=22, color=WHITE)

# 강조 박스
rect(sl, Inches(0.75), Inches(6.25), Inches(11.5), Inches(0.95), CARD)
rect(sl, Inches(0.75), Inches(6.25), Inches(0.05), Inches(0.95), ACCENT)
txt(sl, "피스메이커가 물어볼 것 →  \"그 종목, 왜 사셨어요?\"",
    Inches(1.0), Inches(6.4), Inches(11.0), Inches(0.65),
    size=22, bold=True, color=ACCENT)

# ════════════════════════════════════════════════════════
# SL 07 · 컬럼 분류판
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)

# 좌 컬럼
rect(sl, Inches(0), Inches(0), Inches(6.5), H, RGBColor(0x10,0x10,0x10))
# 우 컬럼
rect(sl, Inches(6.83), Inches(0), Inches(6.5), H, RGBColor(0x0D,0x0D,0x0D))

# 헤더
rect(sl, Inches(0), Inches(0), Inches(6.5), Inches(1.3), CARD2)
rect(sl, Inches(6.83), Inches(0), Inches(6.5), Inches(1.3), CARD)

txt(sl, "조건으로 산 것", Inches(0.2), Inches(0.3), Inches(6.0), Inches(0.7),
    size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "감으로 산 것", Inches(7.0), Inches(0.3), Inches(6.0), Inches(0.7),
    size=28, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)

# 헤더 하단 강조선
rect(sl, Inches(0), Inches(1.3), Inches(6.5), Inches(0.05), ACCENT)
rect(sl, Inches(6.83), Inches(1.3), Inches(6.5), Inches(0.05), GRAY)

# 힌트
txt(sl, "\"이 조건이 충족되면 산다\" 가 있었다면",
    Inches(0.2), Inches(1.45), Inches(6.0), Inches(0.4),
    size=12, color=GRAY, align=PP_ALIGN.CENTER)
txt(sl, "뉴스 · 지인 추천 · 왠지 · 감 · 기분",
    Inches(7.0), Inches(1.45), Inches(6.0), Inches(0.4),
    size=12, color=GRAY, align=PP_ALIGN.CENTER)

# 포스트잇 가이드 칸 (4x2)
for row in range(4):
    for col,base in enumerate([Inches(0.2), Inches(7.0)]):
        t = Inches(2.0) + Inches(1.3)*row
        rect(sl, base, t, Inches(6.0), Inches(1.1), CARD2)
        line(sl, base, t+Inches(1.1)-Inches(0.03),
             Inches(6.0), Inches(0.03), RGBColor(0x2A,0x2A,0x2A))

# 중앙 구분선
rect(sl, Inches(6.5), Inches(0), Inches(0.33), H, BG)
line(sl, Inches(6.63), Inches(0), Inches(0.04), H, CARD2)

txt(sl, "← 포스트잇을 칸에 붙여주세요",
    Inches(2), Inches(7.1), Inches(9), Inches(0.35),
    size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SL 08 · 함께 바라보기
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "잠깐, 다 같이 이걸 한번 봐요.",
    Inches(0.75), Inches(0.55), Inches(12), Inches(1.6),
    size=50, bold=True, color=WHITE)

rect(sl, Inches(0.75), Inches(2.4), Inches(11.5), Inches(0.05), CARD2)

txt(sl, "\"어때요, 보시니까?\"",
    Inches(0.75), Inches(2.65), Inches(11), Inches(0.7),
    size=30, color=LGRAY)
txt(sl, "\"뭔가 느끼시는 거 있어요?\"",
    Inches(0.75), Inches(3.4), Inches(11), Inches(0.7),
    size=30, color=LGRAY)

rect(sl, Inches(0.75), Inches(4.4), Inches(11.5), Inches(1.1), CARD)
rect(sl, Inches(0.75), Inches(4.4), Inches(0.05), Inches(1.1), ACCENT)
txt(sl,
    "이거 제가 만든 게 아니에요.  오늘 여러분이 채운 거예요.",
    Inches(1.0), Inches(4.57), Inches(11.0), Inches(0.7),
    size=22, bold=True, color=ACCENT)

txt(sl, "※ 피스메이커는 분석하지 않습니다. 30초 이상 침묵도 허용.",
    Inches(0.75), Inches(7.05), Inches(10), Inches(0.35),
    size=11, color=GRAY)

# ════════════════════════════════════════════════════════
# SL 09 · 공통점 발견 — 조건 유도 1단계
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "잠깐, 질문 하나 더.",
    Inches(0.75), Inches(0.55), Inches(10), Inches(0.55),
    size=16, color=ACCENT, bold=True)
txt(sl, "\"감으로 산 것\"에서\n공통점이 보이나요?",
    Inches(0.75), Inches(1.2), Inches(10), Inches(2.0),
    size=46, bold=True, color=WHITE)

line(sl, Inches(0.75), Inches(3.35), Inches(11.5), Inches(0.03))

hints = [
    "이유가 외부에 있다  —  뉴스, 지인, 커뮤니티",
    "내가 검증한 게 없다  —  그냥 '오를 것 같았다'",
    "같은 상황에서 다른 날 다른 판단을 했다",
]
for i,h in enumerate(hints):
    t = Inches(3.7) + Inches(0.95)*i
    rect(sl, Inches(0.75), t+Inches(0.12), Inches(0.3), Inches(0.3),
         ACCENT)
    txt(sl, h, Inches(1.2), t, Inches(10.5), Inches(0.6),
        size=20, color=LGRAY)

txt(sl, "그렇다면, 조건으로 산 것과의 차이는 뭘까요?",
    Inches(0.75), Inches(6.75), Inches(11), Inches(0.45),
    size=16, bold=True, color=ACCENT)

# ════════════════════════════════════════════════════════
# SL 10 · 조건 유도 2단계 — 대비
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)

# 좌우 분할
rect(sl, Inches(0), Inches(0), Inches(6.5), H, RGBColor(0x0A,0x0A,0x0A))
rect(sl, Inches(6.83), Inches(0), Inches(6.5), H, RGBColor(0x08,0x12,0x00))

# 헤더
txt(sl, "기준 없이 산 경우", Inches(0.2), Inches(0.5), Inches(6.0), Inches(0.6),
    size=22, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)
txt(sl, "기준을 갖고 산 경우", Inches(7.0), Inches(0.5), Inches(6.0), Inches(0.6),
    size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
line(sl, Inches(0), Inches(1.15), Inches(6.5), Inches(0.04), GRAY)
line(sl, Inches(6.83), Inches(1.15), Inches(6.5), Inches(0.04), ACCENT)

left_items = [
    "왜 샀는지 설명하기 어렵다",
    "손실 나도 왜 잘못됐는지 모른다",
    "다음에 또 같은 실수를 한다",
    "수익이 나도 운인지 실력인지 모른다",
]
right_items = [
    "왜 샀는지 3가지 이상 말할 수 있다",
    "손실 나면 어디서 틀렸는지 안다",
    "같은 실수를 반복하지 않는다",
    "수익을 실력으로 재현할 수 있다",
]
for i,(lft,rgt) in enumerate(zip(left_items,right_items)):
    t = Inches(1.5) + Inches(1.2)*i
    rect(sl, Inches(0.3), t+Inches(0.1), Inches(0.22), Inches(0.22), GRAY)
    txt(sl, lft, Inches(0.7), t, Inches(5.6), Inches(0.65), size=19, color=LGRAY)
    rect(sl, Inches(7.1), t+Inches(0.1), Inches(0.22), Inches(0.22), ACCENT)
    txt(sl, rgt, Inches(7.5), t, Inches(5.6), Inches(0.65), size=19, color=WHITE)

# 구분선
rect(sl, Inches(6.5), Inches(0), Inches(0.33), H, BG)
line(sl, Inches(6.63), Inches(0), Inches(0.04), H, CARD2)

txt(sl, "이 차이를 만드는 게 뭘까요?",
    Inches(0), Inches(6.85), W, Inches(0.45),
    size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SL 11 · 피크 — 빈 질문 (유도)
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)

# 상단 미세 라인
rect(sl, Inches(0), Inches(0), W, Inches(0.05), ACCENT)

txt(sl, "그래서 결국,",
    Inches(0), Inches(1.2), W, Inches(0.8),
    size=28, color=LGRAY, align=PP_ALIGN.CENTER)
txt(sl, "우리에게 필요한 건?",
    Inches(0), Inches(2.1), W, Inches(1.4),
    size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 빈 답변 공간
rect(sl, Inches(3.5), Inches(3.8), Inches(6.3), Inches(1.5), CARD)
rect(sl, Inches(3.5), Inches(3.8), Inches(6.3), Inches(0.05), ACCENT)
txt(sl, "여러분의 언어로 채워주세요.",
    Inches(3.5), Inches(4.3), Inches(6.3), Inches(0.6),
    size=16, color=GRAY, align=PP_ALIGN.CENTER)

# 하단 안내
txt(sl, "※  피스메이커는 이 빈칸을 채우지 않습니다. 피스메이트가 먼저 말할 때까지 기다립니다.",
    Inches(0), Inches(6.9), W, Inches(0.4),
    size=11, color=GRAY, align=PP_ALIGN.CENTER)

rect(sl, Inches(0), H-Inches(0.05), W, Inches(0.05), ACCENT)

# ════════════════════════════════════════════════════════
# SL 12 · 정점 — "조건" 등장 순간
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0), Inches(0), W, Inches(0.05), ACCENT)
rect(sl, Inches(0), H-Inches(0.05), W, Inches(0.05), ACCENT)

txt(sl, "맞습니다.", Inches(0), Inches(0.8), W, Inches(0.7),
    size=24, color=LGRAY, align=PP_ALIGN.CENTER)

txt(sl, "\"  조  건  \"",
    Inches(0), Inches(1.6), W, Inches(2.5),
    size=90, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

line(sl, Inches(2.0), Inches(4.3), Inches(9.3), Inches(0.04))

txt(sl,
    "검증된 조건을  ·  반복 실행할 수 있을 때  ·  수익이 일어난다.",
    Inches(0), Inches(4.55), W, Inches(0.6),
    size=20, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl,
    "이 단어는 피스메이커가 먼저 말하지 않습니다.\n"
    "피스메이트 입에서 먼저 나오는 순간이 오늘 1회차의 정점입니다.",
    Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.9),
    size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SL 13 · 과제 카드
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "다음 주 과제", Inches(0.75), Inches(0.55), Inches(10), Inches(0.65),
    size=40, bold=True, color=WHITE)
line(sl, Inches(0.75), Inches(1.3), Inches(11.5), Inches(0.03))

# 과제 박스
rect(sl, Inches(0.75), Inches(1.6), Inches(11.5), Inches(3.2), CARD)
rect(sl, Inches(0.75), Inches(1.6), Inches(0.05), Inches(3.2), ACCENT)

task_items = [
    ("01", "이번 주 내가 관심 있는 종목 1개를 골라오세요."),
    ("02", "왜 그 종목인지  —  조건 3가지를 직접 적어오세요."),
    ("03", "틀려도 됩니다. 찾는 과정이 중요합니다."),
]
for i,(n,t_) in enumerate(task_items):
    tp = Inches(1.9) + Inches(0.95)*i
    txt(sl, n, Inches(1.0), tp+Inches(0.06), Inches(0.7), Inches(0.45),
        size=13, bold=True, color=ACCENT)
    txt(sl, t_, Inches(1.8), tp, Inches(10.0), Inches(0.6),
        size=21, color=WHITE if i<2 else LGRAY, bold=(i==0))

# 씨앗 질문 박스
rect(sl, Inches(0.75), Inches(5.1), Inches(11.5), Inches(1.4), CARD2)
rect(sl, Inches(0.75), Inches(5.1), Inches(11.5), Inches(0.04), CARD2)
txt(sl, "집에 가시면서 한 번만 생각해보세요.",
    Inches(1.0), Inches(5.25), Inches(11.0), Inches(0.5),
    size=15, bold=True, color=ACCENT)
txt(sl, "나는 몇 번이나 같은 이유로 같은 실수를 반복했을까.",
    Inches(1.0), Inches(5.75), Inches(11.0), Inches(0.55),
    size=22, color=WHITE)
txt(sl, "지금 대답 안 하셔도 됩니다. 그냥 가지고 가세요.",
    Inches(1.0), Inches(6.4), Inches(11.0), Inches(0.45),
    size=14, color=LGRAY)

# ════════════════════════════════════════════════════════
# SL 14 · 클로징
# ════════════════════════════════════════════════════════
sl = blank(); bg(sl)
rect(sl, Inches(0.55), Inches(0.5), Inches(0.05), Inches(6.5), ACCENT)

txt(sl, "오늘 우리가 한 것,", Inches(0.75), Inches(0.55), Inches(11), Inches(0.6),
    size=18, color=LGRAY)
txt(sl, "사실 되게 불편한\n시간이었을 거예요.",
    Inches(0.75), Inches(1.2), Inches(11), Inches(2.0),
    size=48, bold=True, color=WHITE)

line(sl, Inches(0.75), Inches(3.35), Inches(11.5), Inches(0.03))

txt(sl, "내가 감으로 샀다는 걸 인정하는 게 쉽지 않으니까요.",
    Inches(0.75), Inches(3.55), Inches(11), Inches(0.6),
    size=22, color=ACCENT)

txt(sl,
    "다음 주에는 그 조건을 직접 찾아보는 시간인데 —\n"
    "솔직히 말씀드리면, 다음 주도 쉽지 않을 거예요.\n"
    "그걸 경험해야 그다음이 의미 있어집니다.",
    Inches(0.75), Inches(4.3), Inches(11), Inches(1.8),
    size=20, color=LGRAY)

rect(sl, Inches(0.75), Inches(6.5), Inches(11.5), Inches(0.7), CARD)
txt(sl, "오늘 오셔서 감사합니다.  편하게 가세요.",
    Inches(1.0), Inches(6.6), Inches(11.0), Inches(0.5),
    size=20, bold=True, color=WHITE)

# 워터마크
txt(sl, "돈복사 커뮤니티  ·  1회차",
    Inches(0.75), Inches(7.1), Inches(8), Inches(0.3),
    size=11, color=GRAY)

# ── 저장 ────────────────────────────────────────────────
out = r"C:\haessje\라임아카데미\돈복사_1회차.pptx"
prs.save(out)
print("저장 완료:", out)
print("슬라이드 수:", len(prs.slides))
