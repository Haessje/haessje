# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── 컬러 ────────────────────────────────────────
BG     = RGBColor(0x0D,0x0D,0x0D)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LIME   = RGBColor(0xD4,0xFF,0x00)
GRAY   = RGBColor(0x55,0x55,0x55)
LGRAY  = RGBColor(0x99,0x99,0x99)
CARD   = RGBColor(0x1A,0x1A,0x1A)
CARD2  = RGBColor(0x22,0x22,0x22)
RED_C  = RGBColor(0xFF,0x44,0x44)   # 양봉
BLUE_C = RGBColor(0x44,0x88,0xFF)  # 음봉
GREEN  = RGBColor(0x00,0xCC,0x66)
YELL   = RGBColor(0xFF,0xCC,0x00)
ORANGE = RGBColor(0xFF,0x88,0x00)

# ── 공통 헬퍼 ────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(sl, color=BG):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = color

def box(sl, l, t, w, h, color, radius=False):
    s = sl.shapes.add_shape(5 if radius else 1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(sl, s, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b = sl.shapes.add_textbox(l, t, w, h)
    b.word_wrap = wrap
    tf = b.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    r.font.name = "맑은 고딕"
    return b

def txt_multi(sl, lines, l, t, w, h, size=13, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, line_spacing=1.0):
    """여러 줄을 paragraph로 추가"""
    b = sl.shapes.add_textbox(l, t, w, h)
    b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "맑은 고딕"
    return b

def header(sl, title, sub='', accent=LIME):
    box(sl, 0, 0, W, Inches(1.05), accent)
    txt(sl, title, Inches(0.4), Inches(0.1), Inches(10), Inches(0.55),
        size=22, bold=True, color=BG)
    if sub:
        txt(sl, sub, Inches(0.4), Inches(0.62), Inches(11), Inches(0.38),
            size=12, color=BG)

def footer_bar(sl, pg):
    box(sl, 0, H - Inches(0.42), W, Inches(0.42), CARD)
    txt(sl, '돈복사 커뮤니티  ·  1주차 기초 교육자료',
        Inches(0.4), H - Inches(0.38), Inches(8), Inches(0.32),
        size=9, color=LGRAY)
    txt(sl, str(pg), W - Inches(0.8), H - Inches(0.38),
        Inches(0.4), Inches(0.32), size=9, color=LGRAY, align=PP_ALIGN.RIGHT)

def card(sl, l, t, w, h, title, title_color=LIME, title_size=14):
    box(sl, l, t, w, h, CARD2, radius=True)
    box(sl, l, t, w, Inches(0.05), title_color)
    txt(sl, title, l + Inches(0.25), t + Inches(0.1),
        w - Inches(0.3), Inches(0.45), size=title_size, bold=True, color=title_color)
    return l + Inches(0.25), t + Inches(0.6)   # 본문 시작 좌표 반환

def candle_shape(sl, cx, cy_top, body_h, wick_t, wick_b, is_up=True, bw=Inches(0.18)):
    """캔들 한 개 그리기 (cx, cy_top = 몸통 상단 기준)"""
    color = RED_C if is_up else BLUE_C
    wick_w = Inches(0.03)
    # 위 꼬리
    if wick_t > 0:
        s = sl.shapes.add_shape(1, cx - wick_w/2, cy_top - wick_t, wick_w, wick_t)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    # 몸통
    s = sl.shapes.add_shape(1, cx - bw/2, cy_top, bw, body_h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    # 아래 꼬리
    if wick_b > 0:
        s = sl.shapes.add_shape(1, cx - wick_w/2, cy_top + body_h, wick_w, wick_b)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def i(v): return Inches(v)

# ════════════════════════════════════════════════════
# SL 01 · 표지
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
box(sl, 0, 0, W, i(1.1), LIME)
txt(sl, '돈복사 커뮤니티  ·  1주차 기초 교육자료',
    i(0.4), i(0.25), i(10), i(0.55), size=18, bold=True, color=BG)

txt(sl, '주식 기초', i(0.5), i(1.5), i(9), i(1.4), size=72, bold=True, color=WHITE)
txt(sl, '완전 정복', i(0.5), i(2.85), i(9), i(1.4), size=72, bold=True, color=LIME)

txt(sl, '캔들차트  ·  주식 용어  ·  차트 보는 법',
    i(0.5), i(4.35), i(9), i(0.5), size=18, color=LGRAY)
txt(sl, '중고등학생도 이해할 수 있게 쉽게 풀어드립니다.',
    i(0.5), i(4.9), i(9), i(0.45), size=14, color=GRAY)

# 표지 캔들 일러스트
candle_data = [
    (i(10.0), i(2.2), i(0.8), i(0.3), i(0.25), True),
    (i(10.6), i(1.9), i(1.1), i(0.4), i(0.2),  True),
    (i(11.2), i(2.5), i(0.6), i(0.15),i(0.5),  False),
    (i(11.8), i(2.1), i(0.95),i(0.3), i(0.3),  True),
    (i(12.4), i(1.7), i(1.3), i(0.35),i(0.25), True),
]
for cx, cy, bh, wt, wb, up in candle_data:
    candle_shape(sl, cx, cy, bh, wt, wb, up, bw=i(0.28))

box(sl, 0, H - i(0.45), W, i(0.45), CARD)
txt_items = ['양봉/음봉', '시가·고가·저가·종가', '거래량', '이동평균선', '상한가/하한가']
for k, t_ in enumerate(txt_items):
    bx = i(0.4) + i(2.45)*k
    box(sl, bx, H-i(0.38), i(2.2), i(0.3), CARD2, radius=True)
    txt(sl, t_, bx+i(0.1), H-i(0.38), i(2.1), i(0.3),
        size=10, color=LGRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SL 02 · 주식이란?
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
header(sl, '01  주식이란 무엇인가?', '회사의 조각을 사는 것')
footer_bar(sl, 2)

# 피자 비유 박스
box(sl, i(0.4), i(1.2), W - i(0.8), i(1.85), CARD, radius=True)
box(sl, i(0.4), i(1.2), i(0.08), i(1.85), LIME)
txt(sl, '🍕  피자 비유로 이해하기', i(0.65), i(1.3), i(10), i(0.4),
    size=14, bold=True, color=LIME)
txt_multi(sl, [
    '삼성전자라는 피자 한 판이 있다고 생각해봐요.',
    '이 피자를 1억 조각으로 잘랐어요. 그게 바로 "주식"이에요.',
    '내가 그 조각 1개를 사면  →  삼성전자의 아주 작은 주인(주주)이 되는 거예요!',
], i(0.65), i(1.72), W - i(1.1), i(1.1), size=13, color=WHITE)

# 도식 (회사 → 주식 → 투자자)
boxes_info = [
    (i(0.5),  i(3.3), '회사\n(삼성전자)', LIME,  BG,    18),
    (i(5.1),  i(3.3), '주식\n발행',       CARD2, LIME,  18),
    (i(9.7),  i(3.3), '나(투자자)\n주주!', CARD2, WHITE, 18),
]
for bx, by, label, bg_c, fg_c, fs in boxes_info:
    box(sl, bx, by, i(3.4), i(1.7), bg_c, radius=True)
    lines = label.split('\n')
    for li, ln in enumerate(lines):
        txt(sl, ln, bx, by + i(0.35) + i(0.55)*li, i(3.4), i(0.55),
            size=fs, bold=True, color=fg_c, align=PP_ALIGN.CENTER)

# 화살표 대체 (박스)
for ax in [i(4.0), i(8.6)]:
    box(sl, ax, i(4.0), i(1.0), i(0.1), LGRAY)
    box(sl, ax + i(0.6), i(3.85), i(0.4), i(0.4), LGRAY)

# 핵심 포인트
box(sl, i(0.4), i(5.3), W - i(0.8), i(1.75), CARD, radius=True)
txt(sl, '핵심 포인트', i(0.65), i(5.4), i(3), i(0.4),
    size=13, bold=True, color=LIME)
txt_multi(sl, [
    '① 주식을 사면 그 회사의 주인(주주)이 된다',
    '② 회사가 잘되면 주가가 오르고, 못되면 내린다',
    '③ 주가는 매일 수많은 사람들의 매수/매도로 결정된다',
], i(0.65), i(5.82), W - i(1.1), i(1.1), size=13, color=WHITE)

# ════════════════════════════════════════════════════
# SL 03 · 캔들 구조
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
header(sl, '02  캔들차트란?', '하루의 주가 움직임을 초 모양으로 표현한 차트')
footer_bar(sl, 3)

txt(sl, '하나의 캔들 = 하루(또는 1시간, 1분 등)의 가격 정보를 담고 있어요.',
    i(0.4), i(1.15), i(12), i(0.4), size=13, color=LGRAY)

# 큰 캔들 해부도 (양봉)
cx1 = i(2.8)
cy_wt_top = i(1.9)    # 위꼬리 시작
wick_t_h = i(0.8)
body_top = cy_wt_top + wick_t_h   # 몸통 상단
body_h = i(2.4)
wick_b_h = i(0.7)

candle_shape(sl, cx1, body_top, body_h, wick_t_h, wick_b_h, True, bw=i(0.55))

# 라벨 연결선 + 텍스트
label_items = [
    (cx1, cy_wt_top + i(0.05),           '고가  (High)',  '그날 가장 높았던 가격',    WHITE),
    (cx1, body_top + i(0.05),             '종가  (Close)', '장 마감 때 가격 (양봉)',   RED_C),
    (cx1, body_top + body_h/2 - i(0.15), '몸통  (Body)',  '시가와 종가 사이',         LGRAY),
    (cx1, body_top + body_h - i(0.05),   '시가  (Open)',  '장 시작 때 가격 (양봉)',   RED_C),
    (cx1, body_top + body_h + wick_b_h - i(0.05), '저가  (Low)', '그날 가장 낮았던 가격', WHITE),
]
for lx, ly, term, desc, tc in label_items:
    box(sl, lx + i(0.4), ly, i(0.5), i(0.04), LGRAY)
    txt(sl, term, lx + i(0.98), ly - i(0.12), i(2.0), i(0.35),
        size=13, bold=True, color=tc)
    txt(sl, desc, lx + i(0.98), ly + i(0.2), i(3.5), i(0.35),
        size=11, color=LGRAY)

# 음봉 (우측)
cx2 = i(8.5)
candle_shape(sl, cx2, body_top, body_h, wick_t_h, wick_b_h, False, bw=i(0.55))

txt(sl, '음봉 (파란색)', cx2 - i(0.9), body_top - i(0.5), i(2.5), i(0.4),
    size=14, bold=True, color=BLUE_C, align=PP_ALIGN.CENTER)

# 음봉 시가/종가 반대 설명
txt(sl, '↑ 시가 (Open)', cx2 + i(0.42), body_top + i(0.05), i(2.5), i(0.35),
    size=12, color=BLUE_C)
txt(sl, '↑ 종가 (Close)', cx2 + i(0.42), body_top + body_h - i(0.05), i(2.5), i(0.35),
    size=12, color=BLUE_C)
txt(sl, '※ 음봉은 시가가 위, 종가가 아래\n   (종가 < 시가  →  가격이 내린 날)',
    cx2 - i(1.2), body_top + body_h + i(0.6), i(4.0), i(0.8),
    size=11, color=LGRAY)

# 하단 요약 박스
box(sl, i(0.4), i(6.5), i(6.0), i(0.58), RGBColor(0x1A,0x08,0x08), radius=True)
txt(sl, '🔴 양봉  =  종가 > 시가  →  가격이 올랐다',
    i(0.65), i(6.55), i(5.7), i(0.45), size=13, bold=True, color=RED_C)
box(sl, i(6.6), i(6.5), i(6.33), i(0.58), RGBColor(0x08,0x08,0x1A), radius=True)
txt(sl, '🔵 음봉  =  종가 < 시가  →  가격이 내렸다',
    i(6.85), i(6.55), i(6.0), i(0.45), size=13, bold=True, color=BLUE_C)

# ════════════════════════════════════════════════════
# SL 04 · 꼬리의 의미
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
header(sl, '03  캔들 꼬리가 말하는 것', '꼬리 하나에도 시장의 심리가 담겨 있어요')
footer_bar(sl, 4)

tail_cases = [
    {
        'title': '① 윗꼬리가 길다',
        'color': YELL,
        'body_h': i(0.9), 'wt': i(1.5), 'wb': i(0.25),
        'is_up': True,
        'cx': i(1.5), 'cy': i(2.0),
        'desc': [
            '높이 올랐다가 다시 내려왔다는 뜻',
            '→ 고점에서 "팔자"가 많았음',
            '→ 상승 힘이 약해지는 신호일 수 있음',
        ],
    },
    {
        'title': '② 아랫꼬리가 길다',
        'color': GREEN,
        'body_h': i(0.9), 'wt': i(0.25), 'wb': i(1.5),
        'is_up': True,
        'cx': i(5.2), 'cy': i(2.0),
        'desc': [
            '많이 내려갔다가 다시 올라왔다는 뜻',
            '→ 저점에서 "사자"가 강하게 들어옴',
            '→ 하락이 약해지는 신호일 수 있음',
        ],
    },
    {
        'title': '③ 꼬리가 없는 캔들',
        'color': LIME,
        'body_h': i(1.5), 'wt': i(0.0), 'wb': i(0.0),
        'is_up': True,
        'cx': i(8.9), 'cy': i(2.0),
        'desc': [
            '시가=저가, 종가=고가인 경우',
            '→ 하루 종일 한 방향으로만 움직임',
            '→ 매우 강한 추세 신호',
        ],
    },
    {
        'title': '④ 위아래 꼬리 모두 길다',
        'color': ORANGE,
        'body_h': i(0.5), 'wt': i(1.0), 'wb': i(1.0),
        'is_up': False,
        'cx': i(12.0), 'cy': i(2.0),
        'desc': [
            '위아래로 많이 흔들렸다는 뜻',
            '→ 매수/매도 힘겨루기가 치열함',
            '→ 방향성이 불분명한 혼조 상태',
        ],
    },
]

for case in tail_cases:
    cx = case['cx']
    cy = case['cy']
    bh = case['body_h']
    wt = case['wt']
    wb = case['wb']

    candle_shape(sl, cx, cy, bh, wt, wb, case['is_up'], bw=i(0.3))

    txt(sl, case['title'], cx - i(1.5), cy + bh + wb + i(0.15),
        i(3.0), i(0.4), size=12, bold=True, color=case['color'],
        align=PP_ALIGN.CENTER)

    for j, d in enumerate(case['desc']):
        txt(sl, d, cx - i(1.5), cy + bh + wb + i(0.62) + i(0.42)*j,
            i(3.0), i(0.4), size=10, color=LGRAY, align=PP_ALIGN.CENTER)

# 구분선
for xv in [i(3.5), i(7.0), i(10.5)]:
    box(sl, xv, i(1.1), i(0.02), i(6.2), CARD2)

# ════════════════════════════════════════════════════
# SL 05 · 핵심 용어 1 (가격)
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
header(sl, '04  핵심 용어 ①  —  가격 관련', '시가·고가·저가·종가 완벽 이해')
footer_bar(sl, 5)

terms1 = [
    ('시가 (Open)',   '장이 오전 9시에 열릴 때 첫 번째로 거래된 가격',
     '비유: 마라톤의 출발선 가격', '삼성전자 오늘 시가  70,000원', RED_C),
    ('종가 (Close)',  '장이 오후 3시 30분에 닫힐 때 마지막으로 거래된 가격',
     '비유: 마라톤의 결승선 가격', '삼성전자 오늘 종가  72,000원', RED_C),
    ('고가 (High)',   '그날 하루 중 가장 높았던 가격\n잠깐이라도 그 가격에 거래됐으면 고가',
     '', '오늘 고가  74,500원', LIME),
    ('저가 (Low)',    '그날 하루 중 가장 낮았던 가격\n잠깐이라도 그 가격에 거래됐으면 저가',
     '', '오늘 저가  69,800원', BLUE_C),
]

positions = [
    (i(0.4),  i(1.1)),
    (i(6.9),  i(1.1)),
    (i(0.4),  i(4.05)),
    (i(6.9),  i(4.05)),
]
cw, ch = i(6.2), i(2.7)

for (term, meaning, sub, example, ac), (tx, ty) in zip(terms1, positions):
    box(sl, tx, ty, cw, ch, CARD2, radius=True)
    box(sl, tx, ty, cw, i(0.06), ac)
    txt(sl, term, tx + i(0.25), ty + i(0.12), cw - i(0.3), i(0.5),
        size=16, bold=True, color=ac)
    mlines = meaning.split('\n')
    for mi, ml in enumerate(mlines):
        txt(sl, ml, tx + i(0.25), ty + i(0.7) + i(0.5)*mi, cw - i(0.3), i(0.45),
            size=12, color=WHITE)
    if sub:
        txt(sl, sub, tx + i(0.25), ty + i(1.25), cw - i(0.3), i(0.4),
            size=11, color=LGRAY)
    box(sl, tx + i(0.2), ty + ch - i(0.62), cw - i(0.4), i(0.5), CARD)
    txt(sl, '예)  ' + example, tx + i(0.4), ty + ch - i(0.57), cw - i(0.6), i(0.4),
        size=11, color=LGRAY)

# 하단 공식
box(sl, i(0.4), i(7.0), W - i(0.8), i(0.42), CARD)
txt(sl, '공식:  양봉 = 종가 > 시가    /    음봉 = 종가 < 시가    /    고가 ≥ 종가,시가 ≥ 저가',
    i(0.7), i(7.03), W - i(1.0), i(0.35), size=12, bold=True, color=LIME)

# ════════════════════════════════════════════════════
# SL 06 · 핵심 용어 2 (거래)
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
header(sl, '05  핵심 용어 ②  —  매매 & 거래', '매수·매도·거래량·시가총액·상한가·하한가')
footer_bar(sl, 6)

terms2 = [
    ('매수 (BUY)', '주식을 돈을 주고 사는 행위\n비유: 편의점에서 물건을 집어드는 것',
     '10주 매수  →  내 계좌에 주식이 들어온다', GREEN),
    ('매도 (SELL)', '갖고 있던 주식을 파는 행위\n비유: 편의점에 물건을 판매하는 것',
     '10주 매도  →  계좌에 현금이 들어온다', ORANGE),
    ('거래량', '그날 사고 판 주식 수의 총합\n거래량이 많을수록 시장의 관심이 높다',
     '오늘 거래량 1,200만 주  →  관심 많은 종목', LIME),
    ('시가총액', '주가 × 발행 주식 수\n회사 전체가 얼마짜리인지 알 수 있음',
     '주가 70,000원 × 59억 주  =  약 418조 원', YELL),
]

for (term, meaning, example, ac), (tx, ty) in zip(terms2, positions):
    box(sl, tx, ty, cw, ch, CARD2, radius=True)
    box(sl, tx, ty, cw, i(0.06), ac)
    txt(sl, term, tx + i(0.25), ty + i(0.12), cw - i(0.3), i(0.5),
        size=16, bold=True, color=ac)
    mlines = meaning.split('\n')
    for mi, ml in enumerate(mlines):
        txt(sl, ml, tx + i(0.25), ty + i(0.7) + i(0.5)*mi, cw - i(0.3), i(0.45),
            size=12, color=WHITE)
    box(sl, tx + i(0.2), ty + ch - i(0.62), cw - i(0.4), i(0.5), CARD)
    txt(sl, '예)  ' + example, tx + i(0.4), ty + ch - i(0.57), cw - i(0.6), i(0.4),
        size=11, color=LGRAY)

# 상한가/하한가 박스
box(sl, i(0.4), i(7.0), W - i(0.8), i(0.42), CARD)
txt(sl, '상한가 +30%  :  하루에 오를 수 있는 최대폭    |    하한가 -30%  :  하루에 내릴 수 있는 최대폭',
    i(0.7), i(7.03), W - i(1.0), i(0.35), size=12, color=LGRAY)

# ════════════════════════════════════════════════════
# SL 07 · 핵심 용어 3 (차트 분석)
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
header(sl, '06  핵심 용어 ③  —  차트 분석', '이동평균선·지지선·저항선')
footer_bar(sl, 7)

# 이동평균선 설명
box(sl, i(0.4), i(1.1), W - i(0.8), i(1.6), CARD, radius=True)
box(sl, i(0.4), i(1.1), i(0.08), i(1.6), LIME)
txt(sl, '이동평균선 (MA, Moving Average)',
    i(0.65), i(1.18), i(10), i(0.45), size=15, bold=True, color=LIME)
txt_multi(sl, [
    '최근 N일간의 종가 평균을 선으로 이은 것  |  5일선: 단기 추세  /  20일선: 중기 추세',
    '주가가 이동평균선 위  →  상승 추세     /     이동평균선 아래  →  하락 추세',
    '비유: 시험 평균처럼, 최근 점수가 전체 평균보다 높으면 잘하고 있는 것!',
], i(0.65), i(1.65), W - i(1.1), i(1.0), size=12, color=WHITE)

# 미니 차트 (캔들 + 이동평균선 표현)
chart_l, chart_t = i(0.4), i(2.85)
chart_w, chart_h = i(12.53), i(3.2)
box(sl, chart_l, chart_t, chart_w, chart_h, CARD2, radius=True)

candle_info = [
    (0.30, 0.48, True),  (0.48, 0.60, True),  (0.60, 0.50, False),
    (0.50, 0.62, True),  (0.62, 0.72, True),  (0.72, 0.65, False),
    (0.65, 0.78, True),  (0.78, 0.85, True),  (0.85, 0.75, False),
    (0.75, 0.88, True),  (0.88, 0.82, False), (0.82, 0.90, True),
]
n = len(candle_info)
cw_c = (chart_w - i(0.6)) / n
scale = chart_h - i(0.5)
chart_base = chart_t + i(0.25)

for k, (op, cl, up) in enumerate(candle_info):
    cxk = chart_l + i(0.3) + cw_c * k + cw_c/2
    hi = max(op, cl) + 0.04
    lo = min(op, cl) - 0.04
    body_top_y = chart_base + (1-(max(op,cl)))*scale
    bh_k = abs(cl - op) * scale
    wt_k = 0.04 * scale
    wb_k = 0.04 * scale
    candle_shape(sl, cxk, body_top_y, bh_k, wt_k, wb_k, up, bw=cw_c*0.55)

# 이동평균선 대략 표현 (5개 점을 작은 박스로)
closes = [cl for (_, cl, _) in candle_info]
ma5_pts = [sum(closes[max(0,k-4):k+1])/min(k+1,5) for k in range(n)]
ma20_pts = [sum(closes[max(0,k-11):k+1])/min(k+1,12) for k in range(n)]

for k in range(n-1):
    for ma_list, ma_color in [(ma5_pts, YELL), (ma20_pts, RGBColor(0x00,0xAA,0xFF))]:
        y1 = chart_base + (1-ma_list[k])  *scale
        y2 = chart_base + (1-ma_list[k+1])*scale
        cx_a = chart_l + i(0.3) + cw_c*k + cw_c/2
        cx_b = chart_l + i(0.3) + cw_c*(k+1) + cw_c/2
        # 선 근사: 얇은 박스
        mid_x = (cx_a + cx_b)/2
        mid_y = (y1 + y2)/2
        seg_w = cx_b - cx_a
        seg_h = max(abs(y2-y1), i(0.02))
        top_y = min(y1, y2)
        s = sl.shapes.add_shape(1, cx_a, top_y, seg_w, i(0.025))
        s.fill.solid(); s.fill.fore_color.rgb = ma_color
        s.line.fill.background()

# 범례
box(sl, chart_l + i(0.2), chart_t + i(0.15), i(0.6), i(0.08), YELL)
txt(sl, '5일 이동평균선 (단기)', chart_l + i(0.9), chart_t + i(0.1),
    i(3.0), i(0.3), size=10, color=YELL)
box(sl, chart_l + i(4.5), chart_t + i(0.15), i(0.6), i(0.08), RGBColor(0x00,0xAA,0xFF))
txt(sl, '20일 이동평균선 (중기)', chart_l + i(5.2), chart_t + i(0.1),
    i(3.0), i(0.3), size=10, color=RGBColor(0x00,0xAA,0xFF))

# 지지선/저항선
box(sl, i(0.4), i(6.15), i(6.0), i(0.9), CARD, radius=True)
box(sl, i(0.4), i(6.15), i(0.08), i(0.9), GREEN)
txt(sl, '지지선 (Support)', i(0.65), i(6.2), i(4), i(0.4),
    size=13, bold=True, color=GREEN)
txt(sl, '주가가 내려오다 "더 안 내린다"고 버티는 가격대   /   비유: 바닥',
    i(0.65), i(6.6), i(5.5), i(0.38), size=11, color=LGRAY)

box(sl, i(6.9), i(6.15), i(6.03), i(0.9), CARD, radius=True)
box(sl, i(6.9), i(6.15), i(0.08), i(0.9), RED_C)
txt(sl, '저항선 (Resistance)', i(7.15), i(6.2), i(4), i(0.4),
    size=13, bold=True, color=RED_C)
txt(sl, '주가가 올라오다 "더 안 오른다"고 막히는 가격대   /   비유: 천장',
    i(7.15), i(6.6), i(5.5), i(0.38), size=11, color=LGRAY)

# ════════════════════════════════════════════════════
# SL 08 · 실전 차트 읽기 + 퀴즈
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
header(sl, '07  실전!  차트 읽어보기', '지금까지 배운 내용을 한 번에 적용해봐요')
footer_bar(sl, 8)

chart2_l, chart2_t = i(0.4), i(1.1)
chart2_w, chart2_h = i(12.53), i(4.0)
box(sl, chart2_l, chart2_t, chart2_w, chart2_h, CARD2, radius=True)

candle_info2 = [
    (0.40,0.52,True), (0.52,0.60,True), (0.60,0.44,False),
    (0.44,0.50,True), (0.50,0.58,True), (0.58,0.70,True),
    (0.70,0.65,False),(0.65,0.76,True), (0.76,0.84,True),
    (0.84,0.73,False),(0.73,0.88,True), (0.88,0.95,True),
]
n2 = len(candle_info2)
cw2 = (chart2_w - i(0.6)) / n2
scale2 = chart2_h - i(0.5)
base2 = chart2_t + i(0.25)

for k,(op,cl,up) in enumerate(candle_info2):
    cxk = chart2_l + i(0.3) + cw2*k + cw2/2
    body_top_y = base2 + (1-max(op,cl))*scale2
    bh_k = abs(cl-op)*scale2
    wt_k = i(0.12) if k in [2,6,9] else i(0.06)
    wb_k = i(0.12) if k in [2]     else i(0.06)
    candle_shape(sl, cxk, body_top_y, bh_k, wt_k, wb_k, up, bw=cw2*0.6)

# 주석
annots = [
    (2,  '장대음봉!\n주의 구간', RED_C),
    (6,  '음봉 후\n다시 반등', YELL),
    (9,  '조정 후\n강한 상승', GREEN),
]
for idx, note, nc in annots:
    op,cl,up = candle_info2[idx]
    cxk = chart2_l + i(0.3) + cw2*idx + cw2/2
    aty = base2 + (1-max(op,cl))*scale2 - i(0.55)
    for li, ln in enumerate(note.split('\n')):
        txt(sl, ln, cxk - i(0.7), aty + i(0.3)*li, i(1.4), i(0.32),
            size=9, bold=True, color=nc, align=PP_ALIGN.CENTER)

# 날짜 라벨
days = ['1일','2일','3일','4일','5일','6일','7일','8일','9일','10일','11일','12일']
for k,d in enumerate(days):
    cxk = chart2_l + i(0.3) + cw2*k + cw2/2
    txt(sl, d, cxk - i(0.35), chart2_t + chart2_h + i(0.05), i(0.7), i(0.28),
        size=8, color=GRAY, align=PP_ALIGN.CENTER)

# 퀴즈
box(sl, i(0.4), i(5.45), W - i(0.8), i(1.62), CARD, radius=True)
txt(sl, '✏️  연습 퀴즈  —  아래 질문에 답해보세요!',
    i(0.65), i(5.52), i(11), i(0.4), size=13, bold=True, color=LIME)
quizzes = [
    'Q1.  위 차트에서 양봉은 몇 개,  음봉은 몇 개인가요?',
    'Q2.  3일째에 큰 음봉이 나왔을 때 시장 분위기는 어땠을까요?',
    'Q3.  10일째에 조정 후 11~12일 연속 양봉이 나왔다면 어떤 신호일까요?',
]
for qi, q in enumerate(quizzes):
    txt(sl, q, i(0.65), i(6.0) + i(0.35)*qi, W - i(1.1), i(0.32),
        size=12, color=WHITE)

# ════════════════════════════════════════════════════
# SL 09 · 핵심 정리
# ════════════════════════════════════════════════════
sl = blank(); bg(sl)
header(sl, '08  오늘의 핵심 정리', '이것만 알면 차트를 읽을 수 있어요')
footer_bar(sl, 9)

checks = [
    ('주식 = 회사의 조각.  주주가 되면 회사 이익을 나눌 수 있다', True),
    ('캔들 하나 = 하루의 가격 정보  (시가 · 고가 · 저가 · 종가)', True),
    ('빨간 캔들(양봉) = 가격 올랐다     /     파란 캔들(음봉) = 가격 내렸다', True),
    ('긴 윗꼬리 → 고점 매도 강함    /    긴 아랫꼬리 → 저점 매수 강함', True),
    ('거래량이 많을수록 그날 시장의 관심이 높다는 의미', True),
    ('이동평균선 : 최근 N일 평균 가격선  (5일선·20일선 등)', True),
    ('지지선 = 잘 안 내려가는 가격대    /    저항선 = 잘 안 올라가는 가격대', True),
    ('상한가 +30%  /  하한가 -30%  :  하루 최대 등락폭 제한', True),
]

for ki, (text, done) in enumerate(checks):
    row = ki % 4
    col = ki // 4
    tx = i(0.4) + i(6.5)*col
    ty = i(1.15) + i(1.35)*row
    bw_k = i(6.1)
    bh_k = i(1.2)
    box(sl, tx, ty, bw_k, bh_k, CARD2, radius=True)
    box(sl, tx, ty, bw_k, i(0.06), LIME)
    # 체크 원
    box(sl, tx + i(0.15), ty + i(0.3), i(0.45), i(0.45), LIME, radius=True)
    txt(sl, '✓', tx + i(0.17), ty + i(0.28), i(0.42), i(0.42),
        size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, text, tx + i(0.72), ty + i(0.3), bw_k - i(0.85), i(0.6),
        size=11, color=WHITE)

# 마무리 메시지
box(sl, i(0.4), i(6.95), W - i(0.8), i(0.5), LIME, radius=True)
txt(sl, '다음 주부터는 이 지식으로 직접 "좋은 종목의 조건"을 찾아봅니다.  함께 찾아봐요! — 피스메이커',
    i(0.7), i(6.99), W - i(1.2), i(0.4), size=13, bold=True, color=BG)

out = r'C:\haessje\라임아카데미\돈복사_1주차_기초교육.pptx'
prs.save(out)
print('완료:', out, '  슬라이드 수:', len(prs.slides))
