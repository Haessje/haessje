import io
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG     = RGBColor(0x0D, 0x0D, 0x0D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIME   = RGBColor(0xD4, 0xFF, 0x00)
GRAY   = RGBColor(0x88, 0x88, 0x88)
RED    = RGBColor(0xFF, 0x44, 0x44)
BLUE   = RGBColor(0x44, 0x88, 0xFF)
GREEN  = RGBColor(0x00, 0xCC, 0x66)
YELLOW = RGBColor(0xFF, 0xCC, 0x00)
ORANGE = RGBColor(0xFF, 0x99, 0x00)

W = Inches(13.33)
H = Inches(7.5)

plt.rcParams['font.family'] = 'Malgun Gothic'
plt.rcParams['axes.unicode_minus'] = False

def hr(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def txt(slide, text, left, top, width, height,
        size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.color.rgb = color
    r.font.name      = 'Malgun Gothic'

def box(slide, left, top, width, height, fill):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()

def add_img(slide, fig, left, top, width, height):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150,
                bbox_inches='tight', facecolor=fig.get_facecolor())
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width, height)
    plt.close(fig)

def make_ohlc(n=100, seed=7, trend=0.12, vol=1.6):
    np.random.seed(seed)
    r = np.random.normal(trend/100, vol/100, n)
    C = 100 * np.exp(np.cumsum(r))
    O = np.empty(n); O[0] = C[0]
    for i in range(1, n): O[i] = C[i-1]
    wick = np.abs(C) * 0.012
    rng  = np.random.default_rng(seed)
    H = np.maximum(O, C) + np.abs(rng.normal(0, wick, n))
    L = np.minimum(O, C) - np.abs(rng.normal(0, wick, n))
    return O, H, L, C

def draw_candles(ax, O, H, L, C, width=0.6):
    x    = np.arange(len(C))
    bull = C >= O; bear = ~bull
    RD, BL = '#FF4444', '#4488FF'
    for mask, c in [(bull, RD), (bear, BL)]:
        if mask.any():
            ht = np.maximum(np.abs(C[mask]-O[mask]), np.abs(C[mask])*0.001)
            ax.bar(x[mask], ht, bottom=np.minimum(O[mask],C[mask]),
                   width=width, color=c, zorder=4, linewidth=0)
            ax.vlines(x[mask], L[mask], H[mask], color=c, lw=0.7, zorder=3)

def ax_style(ax):
    ax.set_facecolor('#0D0D0D')
    for s in ax.spines.values(): s.set_color('#333333')
    ax.tick_params(colors='#555555', labelsize=8)

# ─────────────────────────────────────────────────────────
# Slide 1 — 타이틀
# ─────────────────────────────────────────────────────────
def slide_title(prs):
    sl = blank_slide(prs)
    box(sl, Inches(0), Inches(0), Inches(0.15), H, LIME)
    txt(sl, '3주차', Inches(0.4), Inches(1.4), Inches(4), Inches(0.7),
        size=18, bold=True, color=LIME)
    txt(sl, '좋은 타점을\n어떻게 잡는가?',
        Inches(0.4), Inches(2.0), Inches(9), Inches(2.8),
        size=52, bold=True, color=WHITE)
    txt(sl, '이동평균선 + 볼린저밴드 — 타점의 신뢰도를 높이는 법',
        Inches(0.4), Inches(4.7), Inches(11), Inches(0.6),
        size=22, color=GRAY)
    txt(sl, '돈복사', Inches(10), Inches(6.7), Inches(3), Inches(0.5),
        size=15, color=GRAY, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────
# Slide 2 — 오늘의 흐름
# ─────────────────────────────────────────────────────────
def slide_agenda(prs):
    sl = blank_slide(prs)
    txt(sl, '오늘의 흐름', Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
        size=14, color=GRAY)
    txt(sl, '2시간, 이렇게 씁니다',
        Inches(0.5), Inches(0.85), Inches(11), Inches(0.7),
        size=32, bold=True, color=WHITE)

    rows = [
        ('0~10분',   '오프닝',    '과제 종목 공유 — 타점 후보 꺼내기',          GRAY,   '#222222'),
        ('10~40분',  '교육 ①',   '이동평균선 — 정배열/역배열로 타점 판단',      LIME,   '#0D1A00'),
        ('40~55분',  '짝활동 ①', '가져온 종목, 이평선으로 체크',               YELLOW, '#1A1500'),
        ('55~85분',  '교육 ②',   '볼린저밴드 — 타점의 신뢰도 높이기',           LIME,   '#0D1A00'),
        ('85~100분', '짝활동 ②', '두 조건 동시 충족하는 타점 찾기',             YELLOW, '#1A1500'),
        ('100~120분','마무리',    '짝별 공유 + 4주차 예고 (검색기 만들기)',       GRAY,   '#222222'),
    ]

    y = Inches(1.85)
    for time, label, desc, col, bg in rows:
        box(sl, Inches(0.5), y, Inches(12.3), Inches(0.65), hr(bg))
        box(sl, Inches(0.5), y, Inches(0.07), Inches(0.65), col)
        txt(sl, time,  Inches(0.7), y+Inches(0.1), Inches(1.4), Inches(0.5), size=13, color=GRAY)
        txt(sl, label, Inches(2.2), y+Inches(0.1), Inches(1.8), Inches(0.5), size=15, bold=True, color=col)
        txt(sl, desc,  Inches(4.1), y+Inches(0.1), Inches(8.8), Inches(0.5), size=15, color=WHITE)
        y += Inches(0.8)

# ─────────────────────────────────────────────────────────
# Slide 3 — 오프닝
# ─────────────────────────────────────────────────────────
def slide_opening(prs):
    sl = blank_slide(prs)
    txt(sl, '오프닝  |  0~10분', Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
        size=14, color=GRAY)
    txt(sl, '지난 주 과제 공유', Inches(0.5), Inches(0.85), Inches(10), Inches(0.65),
        size=36, bold=True, color=WHITE)
    txt(sl, '타점 후보를 찾은 종목을 꺼내주세요',
        Inches(0.5), Inches(1.65), Inches(12), Inches(0.55),
        size=20, color=LIME)

    # 발표 카드 3개
    cards = [
        ('종목명', '어떤 타점 후보를 발견했나요?', '①'),
        ('근거',   '거래량급증 / 눌림목 / 전고점 중 어느 것?', '②'),
        ('판단',   '지금 사도 되겠다고 생각했나요?', '③'),
    ]
    x = Inches(0.5)
    for num, title, desc in cards:
        box(sl, x, Inches(2.5), Inches(3.8), Inches(3.8), hr('#1A1A1A'))
        box(sl, x, Inches(2.5), Inches(3.8), Inches(0.07), LIME)
        txt(sl, num,   x+Inches(0.2), Inches(2.6), Inches(3.4), Inches(0.55),
            size=13, color=GRAY)
        txt(sl, title, x+Inches(0.2), Inches(3.1), Inches(3.5), Inches(0.55),
            size=20, bold=True, color=WHITE)
        txt(sl, desc,  x+Inches(0.2), Inches(3.65), Inches(3.5), Inches(1.4),
            size=15, color=GRAY)
        x += Inches(4.2)

    txt(sl, '각자 30초씩  —  정답 없음, 솔직하게',
        Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
        size=15, color=GRAY)

# ─────────────────────────────────────────────────────────
# Slide 4 — 교육① 이동평균선 개념
# ─────────────────────────────────────────────────────────
def slide_ma_intro(prs):
    sl = blank_slide(prs)
    txt(sl, '교육 ①  |  10~40분', Inches(0.5), Inches(0.3), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '이동평균선 — 추세를 읽는 기준선',
        Inches(0.5), Inches(0.85), Inches(12), Inches(0.65),
        size=32, bold=True, color=WHITE)

    items = [
        ('20일선', '단기 추세  —  주가가 이선 위에 있으면 단기 상승 중', '#D4FF00'),
        ('60일선', '중기 추세  —  2주차에 배운 타점의 신뢰도 확인',     '#FF9900'),
        ('120일선','장기 추세  —  이선 위에 있으면 큰 흐름이 살아있음',  '#4488FF'),
    ]
    y = Inches(1.9)
    for label, desc, col in items:
        box(sl, Inches(0.5), y, Inches(12.3), Inches(0.9), hr('#1A1A1A'))
        box(sl, Inches(0.5), y, Inches(0.07), Inches(0.9), hr(col))
        txt(sl, label, Inches(0.75), y+Inches(0.17), Inches(1.5), Inches(0.55),
            size=20, bold=True, color=hr(col))
        txt(sl, desc,  Inches(2.4),  y+Inches(0.17), Inches(10.2), Inches(0.55),
            size=18, color=WHITE)
        y += Inches(1.05)

    # 핵심 박스
    box(sl, Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.7), hr('#0D1A00'))
    box(sl, Inches(0.5), Inches(5.15), Inches(0.07), Inches(1.7), LIME)
    txt(sl, '정배열  vs  역배열',
        Inches(0.75), Inches(5.25), Inches(10), Inches(0.55),
        size=22, bold=True, color=LIME)
    txt(sl, '정배열 (20 > 60 > 120)  →  주가가 올라가기 좋은 환경\n역배열 (20 < 60 < 120)  →  반등이 와도 저항이 많은 환경',
        Inches(0.75), Inches(5.8), Inches(12), Inches(0.85),
        size=17, color=WHITE)

# ─────────────────────────────────────────────────────────
# Slide 5 — 이동평균선 차트
# ─────────────────────────────────────────────────────────
def slide_ma_chart(prs):
    sl = blank_slide(prs)
    txt(sl, '정배열 — 타점이 살아있는 환경',
        Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
        size=28, bold=True, color=WHITE)

    O, H, L, C = make_ohlc(n=120, seed=3, trend=0.15, vol=1.5)
    n = len(C)
    x = np.arange(n)

    ma20  = np.array([np.mean(C[max(0,i-19):i+1]) for i in range(n)])
    ma60  = np.array([np.mean(C[max(0,i-59):i+1]) for i in range(n)])
    ma120 = np.array([np.mean(C[max(0,i-119):i+1]) for i in range(n)])

    fig, ax = plt.subplots(figsize=(12, 5.5), facecolor='#0D0D0D')
    ax_style(ax)

    draw_candles(ax, O, H, L, C)
    ax.plot(x, ma20,  color='#D4FF00', lw=1.4, label='20일선', zorder=5)
    ax.plot(x, ma60,  color='#FF9900', lw=1.4, label='60일선', zorder=5)
    ax.plot(x, ma120, color='#4488FF', lw=1.4, label='120일선', zorder=5)

    # 정배열 구간 표시
    align_start = 60
    ax.axvspan(align_start, n-1, alpha=0.08, color='#D4FF00')
    ax.annotate('정배열 구간\n(20 > 60 > 120)',
                xy=(75, ma20[75]+1), xytext=(78, ma20[75]+8),
                fontsize=10, color='#D4FF00', fontfamily='Malgun Gothic',
                arrowprops=dict(arrowstyle='->', color='#D4FF00', lw=1.2))

    # 눌림목 타점 표시
    idx = 85
    ax.annotate('눌림목 타점\n(정배열 + 20일선 지지)',
                xy=(idx, L[idx]-0.5),
                xytext=(idx-18, L[idx]-9),
                fontsize=9.5, color='#FF9900', fontfamily='Malgun Gothic',
                arrowprops=dict(arrowstyle='->', color='#FF9900', lw=1.2))

    legend = ax.legend(loc='upper left', fontsize=10,
                       facecolor='#1A1A1A', edgecolor='#333333', labelcolor='white')
    for text in legend.get_texts():
        text.set_fontfamily('Malgun Gothic')

    ax.set_xlim(-1, n+1)
    ax.set_ylabel('가격', color='#666666', fontsize=9, fontfamily='Malgun Gothic')
    plt.tight_layout(pad=0.3)
    add_img(sl, fig, Inches(0.3), Inches(0.95), Inches(12.7), Inches(6.3))

# ─────────────────────────────────────────────────────────
# Slide 6 — 짝활동①
# ─────────────────────────────────────────────────────────
def slide_pair1(prs):
    sl = blank_slide(prs)
    box(sl, Inches(0), Inches(0), W, Inches(0.08), YELLOW)
    txt(sl, '짝활동 ①  |  40~55분',
        Inches(0.5), Inches(0.25), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '내 종목, 이평선으로 체크',
        Inches(0.5), Inches(0.8), Inches(12), Inches(0.75),
        size=34, bold=True, color=WHITE)

    box(sl, Inches(0.5), Inches(1.75), Inches(12.3), Inches(1.0), hr('#1A1500'))
    txt(sl, '미션  —  과제 종목 차트를 열고 짝이랑 이동평균선 상태 확인하기',
        Inches(0.7), Inches(1.85), Inches(11.8), Inches(0.7),
        size=20, bold=True, color=YELLOW)

    steps = [
        ('①', '차트에서 20일선 / 60일선 / 120일선 켜기'),
        ('②', '정배열인가, 역배열인가? 짝이랑 판단하기'),
        ('③', '주가가 이평선 위에 있는가, 아래에 있는가?'),
        ('④', '타점 후보가 이평선 지지를 받는 자리인가?'),
        ('⑤', '결과를 한 문장으로 정리  ("이 종목은 지금 ___ 상태다")'),
    ]
    y = Inches(2.95)
    for num, s in steps:
        txt(sl, num, Inches(0.7), y, Inches(0.5), Inches(0.55),
            size=18, bold=True, color=YELLOW)
        txt(sl, s, Inches(1.2), y, Inches(11.3), Inches(0.55), size=20, color=WHITE)
        y += Inches(0.65)

    txt(sl, '⏱  15분', Inches(10.2), Inches(6.75), Inches(2.6), Inches(0.5),
        size=22, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────
# Slide 7 — 교육② 볼린저밴드 개념
# ─────────────────────────────────────────────────────────
def slide_bb_intro(prs):
    sl = blank_slide(prs)
    txt(sl, '교육 ②  |  55~85분',
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '볼린저밴드 — 타점의 신뢰도를 높인다',
        Inches(0.5), Inches(0.85), Inches(12), Inches(0.65),
        size=32, bold=True, color=WHITE)

    parts = [
        ('상단밴드', '주가가 여기 닿으면 과열 신호  —  단기 조정 가능성',    '#FF4444', '#2A0D0D'),
        ('중심선',  '20일 이동평균  —  이걸 기준으로 위/아래를 판단',        '#D4FF00', '#1A1A00'),
        ('하단밴드', '주가가 여기 닿으면 과매도 신호  —  반등 타점 후보',     '#4488FF', '#0D0D2A'),
    ]
    y = Inches(1.9)
    for label, desc, col, bg in parts:
        box(sl, Inches(0.5), y, Inches(12.3), Inches(0.95), hr(bg))
        box(sl, Inches(0.5), y, Inches(0.07), Inches(0.95), hr(col))
        txt(sl, label, Inches(0.75), y+Inches(0.2), Inches(2.0), Inches(0.55),
            size=20, bold=True, color=hr(col))
        txt(sl, desc,  Inches(2.9),  y+Inches(0.2), Inches(9.8), Inches(0.55),
            size=18, color=WHITE)
        y += Inches(1.1)

    box(sl, Inches(0.5), Inches(5.25), Inches(12.3), Inches(1.8), hr('#0D1A00'))
    box(sl, Inches(0.5), Inches(5.25), Inches(0.07), Inches(1.8), LIME)
    txt(sl, '핵심 조합  —  이평선 + 볼린저밴드',
        Inches(0.75), Inches(5.35), Inches(11), Inches(0.55),
        size=20, bold=True, color=LIME)
    txt(sl, '정배열 상태에서  +  하단밴드 터치 후 반등',
        Inches(0.75), Inches(5.9), Inches(11), Inches(0.55),
        size=20, bold=True, color=WHITE)
    txt(sl, '조건 하나보다 두 가지가 동시에 맞을 때 신뢰도가 올라간다',
        Inches(0.75), Inches(6.4), Inches(11), Inches(0.5),
        size=16, color=GRAY)

# ─────────────────────────────────────────────────────────
# Slide 8 — 볼린저밴드 차트
# ─────────────────────────────────────────────────────────
def slide_bb_chart(prs):
    sl = blank_slide(prs)
    txt(sl, '볼린저밴드 하단 터치 + 정배열 — 신뢰도 높은 타점',
        Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
        size=26, bold=True, color=WHITE)

    O, H, L, C = make_ohlc(n=120, seed=11, trend=0.12, vol=1.8)
    n = len(C)
    x = np.arange(n)

    ma20 = np.array([np.mean(C[max(0,i-19):i+1]) for i in range(n)])
    std  = np.array([np.std( C[max(0,i-19):i+1]) for i in range(n)])
    upper = ma20 + 2 * std
    lower = ma20 - 2 * std

    ma60  = np.array([np.mean(C[max(0,i-59):i+1])  for i in range(n)])
    ma120 = np.array([np.mean(C[max(0,i-119):i+1]) for i in range(n)])

    fig, ax = plt.subplots(figsize=(12, 5.5), facecolor='#0D0D0D')
    ax_style(ax)

    draw_candles(ax, O, H, L, C)

    ax.plot(x, upper, color='#FF4444', lw=1.1, linestyle='--', label='상단밴드', alpha=0.8)
    ax.plot(x, ma20,  color='#D4FF00', lw=1.4, label='중심선(20일)', zorder=5)
    ax.plot(x, lower, color='#4488FF', lw=1.1, linestyle='--', label='하단밴드', alpha=0.8)
    ax.fill_between(x, upper, lower, alpha=0.04, color='#888888')

    ax.plot(x, ma60,  color='#FF9900', lw=1.0, label='60일선', alpha=0.7)
    ax.plot(x, ma120, color='#9966FF', lw=1.0, label='120일선', alpha=0.7)

    # 하단밴드 터치 타점 찾기
    touch_pts = [i for i in range(20, n-5)
                 if L[i] <= lower[i] and C[i] > lower[i] and C[i] > O[i]]
    if touch_pts:
        for tp in touch_pts[:2]:
            ax.annotate('하단터치\n+ 양봉 반등',
                        xy=(tp, L[tp]-0.5),
                        xytext=(tp+5, L[tp]-7),
                        fontsize=9, color='#4488FF', fontfamily='Malgun Gothic',
                        arrowprops=dict(arrowstyle='->', color='#4488FF', lw=1.2))

    legend = ax.legend(loc='upper left', fontsize=9,
                       facecolor='#1A1A1A', edgecolor='#333333', labelcolor='white',
                       ncol=2)
    for t in legend.get_texts(): t.set_fontfamily('Malgun Gothic')

    ax.set_xlim(-1, n+1)
    ax.set_ylabel('가격', color='#666666', fontsize=9, fontfamily='Malgun Gothic')
    plt.tight_layout(pad=0.3)
    add_img(sl, fig, Inches(0.3), Inches(0.95), Inches(12.7), Inches(6.3))

# ─────────────────────────────────────────────────────────
# Slide 9 — 짝활동②
# ─────────────────────────────────────────────────────────
def slide_pair2(prs):
    sl = blank_slide(prs)
    box(sl, Inches(0), Inches(0), W, Inches(0.08), YELLOW)
    txt(sl, '짝활동 ②  |  85~100분',
        Inches(0.5), Inches(0.25), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '두 조건 동시에 충족하는 타점 찾기',
        Inches(0.5), Inches(0.8), Inches(12), Inches(0.75),
        size=34, bold=True, color=WHITE)

    # 두 조건 박스
    box(sl, Inches(0.5), Inches(1.75), Inches(5.8), Inches(1.3), hr('#0D1A00'))
    txt(sl, '조건 A  —  이동평균선',
        Inches(0.7), Inches(1.85), Inches(5.3), Inches(0.5),
        size=18, bold=True, color=LIME)
    txt(sl, '정배열 상태\n+ 주가가 이평선 위',
        Inches(0.7), Inches(2.3), Inches(5.3), Inches(0.6),
        size=16, color=WHITE)

    txt(sl, '+', Inches(6.4), Inches(2.1), Inches(0.6), Inches(0.7),
        size=32, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    box(sl, Inches(7.1), Inches(1.75), Inches(5.7), Inches(1.3), hr('#0D0D2A'))
    txt(sl, '조건 B  —  볼린저밴드',
        Inches(7.3), Inches(1.85), Inches(5.2), Inches(0.5),
        size=18, bold=True, color=hr('#4488FF'))
    txt(sl, '하단밴드 터치 후 반등\n+ 거래량 동반',
        Inches(7.3), Inches(2.3), Inches(5.2), Inches(0.6),
        size=16, color=WHITE)

    steps = [
        ('①', '차트에 이평선 3개 + 볼린저밴드 동시에 켜기'),
        ('②', '조건 A 확인  (정배열인가?)'),
        ('③', '조건 B 확인  (하단밴드 근처에서 반등 신호 있는가?)'),
        ('④', '두 조건이 동시에 보이는 구간 찾기'),
        ('⑤', '짝에게 "이 자리가 타점이라고 생각한 이유" 설명'),
    ]
    y = Inches(3.3)
    for num, s in steps:
        txt(sl, num, Inches(0.7), y, Inches(0.5), Inches(0.55),
            size=18, bold=True, color=YELLOW)
        txt(sl, s, Inches(1.2), y, Inches(11.3), Inches(0.55), size=20, color=WHITE)
        y += Inches(0.63)

    txt(sl, '⏱  15분', Inches(10.2), Inches(6.75), Inches(2.6), Inches(0.5),
        size=22, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────
# Slide 10 — 핵심 메시지
# ─────────────────────────────────────────────────────────
def slide_key_message(prs):
    sl = blank_slide(prs)
    txt(sl, '오늘의 핵심',
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
        size=14, color=GRAY)

    txt(sl, '조건이 하나일 때보다\n두 개가 겹칠 때 신뢰도가 올라간다',
        Inches(0.5), Inches(1.0), Inches(12.3), Inches(2.2),
        size=40, bold=True, color=WHITE)

    cols = [
        ('조건 1개', '타점 후보', '신뢰도 보통', '#888888', '#1A1A1A'),
        ('조건 2개\n동시 충족', '타점 강화', '신뢰도 높음', '#D4FF00', '#0D1A00'),
        ('조건 3개+\n동시 충족', '최적 타점', '신뢰도 매우 높음', '#00CC66', '#001A0D'),
    ]
    x = Inches(0.5)
    for top, mid, bot, col, bg in cols:
        box(sl, x, Inches(3.5), Inches(3.8), Inches(3.3), hr(bg))
        box(sl, x, Inches(3.5), Inches(3.8), Inches(0.07), hr(col))
        txt(sl, top, x+Inches(0.2), Inches(3.6),  Inches(3.4), Inches(0.9),
            size=20, bold=True, color=hr(col))
        txt(sl, mid, x+Inches(0.2), Inches(4.55), Inches(3.4), Inches(0.55),
            size=22, bold=True, color=WHITE)
        txt(sl, bot, x+Inches(0.2), Inches(5.1),  Inches(3.4), Inches(0.55),
            size=16, color=GRAY)
        x += Inches(4.2)

    txt(sl, '4주차에서는 이 조건들을 자동으로 찾아주는 검색기를 직접 만든다',
        Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.5),
        size=16, bold=True, color=LIME)

# ─────────────────────────────────────────────────────────
# Slide 11 — 마무리
# ─────────────────────────────────────────────────────────
def slide_closing(prs):
    sl = blank_slide(prs)
    txt(sl, '마무리  |  100~120분',
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '오늘 우리가 한 것',
        Inches(0.5), Inches(0.85), Inches(11), Inches(0.65),
        size=36, bold=True, color=WHITE)

    summary = [
        ('이동평균선으로',   '타점이 살아있는 환경인지 판단했다'),
        ('볼린저밴드로',     '하단 터치 후 반등 자리를 확인했다'),
        ('두 조건을 겹쳐서', '신뢰도 높은 타점을 직접 찾아봤다'),
    ]
    y = Inches(1.9)
    for label, val in summary:
        txt(sl, label, Inches(0.5), y, Inches(5.5), Inches(0.55), size=18, color=GRAY)
        txt(sl, val,   Inches(6.2), y, Inches(6.8), Inches(0.55), size=18, bold=True, color=LIME)
        y += Inches(0.75)

    box(sl, Inches(0.5), Inches(3.9), Inches(12.3), Inches(1.05), hr('#0D1A00'))
    txt(sl, '과제  —  오늘 찾은 타점 종목 1개, 진입 이유 2가지 적어오기',
        Inches(0.7), Inches(4.0), Inches(11.8), Inches(0.7),
        size=19, bold=True, color=LIME)
    txt(sl, '(이평선 근거 하나 + 볼린저밴드 근거 하나)',
        Inches(0.7), Inches(4.6), Inches(11.8), Inches(0.4),
        size=15, color=GRAY)

    box(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6), hr('#111111'))
    txt(sl, '4주차 예고  —  검색기 만들기',
        Inches(0.7), Inches(5.32), Inches(11.5), Inches(0.55),
        size=20, bold=True, color=WHITE)
    txt(sl, '오늘 배운 조건들을 키움 / HTS 검색기에 직접 입력\n"내 조건을 자동으로 찾아주는 도구"를 손으로 만든다',
        Inches(0.7), Inches(5.85), Inches(11.5), Inches(0.75),
        size=16, color=GRAY)

    txt(sl, '피스메이커 자가점검  |  차이보존 · 공감≠동감 · 최소한의힘',
        Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
        size=12, color=hr('#444444'))

# ─────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────
prs = new_prs()

slide_title(prs)
slide_agenda(prs)
slide_opening(prs)
slide_ma_intro(prs)
slide_ma_chart(prs)
slide_pair1(prs)
slide_bb_intro(prs)
slide_bb_chart(prs)
slide_pair2(prs)
slide_key_message(prs)
slide_closing(prs)

out = r'C:\haessje\라임아카데미\돈복사_3주차.pptx'
prs.save(out)
print(f'완료: {out}  ({len(prs.slides)}슬라이드)')
