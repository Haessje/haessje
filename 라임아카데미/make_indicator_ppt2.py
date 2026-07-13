"""돈복사 지표기초 PPT v2 - matplotlib 차트 이미지 임베드"""
import io, numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.lines import Line2D
from matplotlib.patches import Patch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

plt.rcParams['font.family'] = 'Malgun Gothic'
plt.rcParams['axes.unicode_minus'] = False

# hex colors for matplotlib
BH='#0D0D0D'; CB='#1A1A1A'; DG='#2A2A2A'
WH='#FFFFFF'; GY='#888888'
AC='#D4FF00'; RD='#FF4444'; BL='#4488FF'
OR='#FFAA00'; GN='#00CC66'

# pptx colors
pBG=RGBColor(0x0D,0x0D,0x0D); pWH=RGBColor(0xFF,0xFF,0xFF)
pAC=RGBColor(0xD4,0xFF,0x00); pRD=RGBColor(0xFF,0x44,0x44)
pBL=RGBColor(0x44,0x88,0xFF); pGY=RGBColor(0x77,0x77,0x77)
pD2=RGBColor(0x22,0x22,0x22); pOR=RGBColor(0xFF,0xAA,0x00)
pGN=RGBColor(0x00,0xCC,0x66)

W=Inches(13.33); H=Inches(7.5)

# ── data helpers ──────────────────────────────────────────────
def sma(d,p):
    r=np.full(len(d),np.nan)
    for i in range(p-1,len(d)): r[i]=np.mean(d[i-p+1:i+1])
    return r

def std_d(d,p):
    r=np.full(len(d),np.nan)
    for i in range(p-1,len(d)): r[i]=np.std(d[i-p+1:i+1])
    return r

def trend(n,s,tr,v,seed):
    np.random.seed(seed)
    return s*np.exp(np.cumsum(np.random.normal(tr/100,v/100,n)))

def ranging(n,ctr,amp,seed):
    np.random.seed(seed); t=np.linspace(0,4*np.pi,n)
    return ctr+amp*np.sin(t)+np.random.normal(0,1.2,n)

# ── chart style ───────────────────────────────────────────────
def ax_s(ax,title='',tc=WH):
    ax.set_facecolor(CB)
    for sp in ax.spines.values(): sp.set_color(DG)
    ax.tick_params(colors=GY,labelsize=8)
    ax.set_xticklabels([]); ax.set_ylabel('')
    ax.grid(True,color=DG,lw=0.5,alpha=0.9,ls='--')
    if title: ax.set_title(title,color=tc,fontsize=11,pad=6,fontweight='bold')

def leg(ax,handles=None,loc='upper left'):
    kw=dict(facecolor='#111111',edgecolor=DG,labelcolor=WH,fontsize=9,framealpha=0.92)
    if handles: ax.legend(handles=handles,loc=loc,**kw)
    else: ax.legend(loc=loc,**kw)

def save(fig):
    buf=io.BytesIO()
    fig.savefig(buf,format='png',dpi=150,bbox_inches='tight',facecolor=BH,edgecolor='none')
    buf.seek(0); plt.close(fig); return buf

# ═══════════════════════════════════════════════════════
# CHART FUNCTIONS
# ═══════════════════════════════════════════════════════

def ch_ma():
    p=trend(120,100,0.18,1.8,42)
    fig,ax=plt.subplots(figsize=(7,4.5)); fig.patch.set_facecolor(BH); ax_s(ax)
    x=np.arange(len(p))
    ax.plot(x,p,       color=WH,lw=1.2,alpha=0.65,label='주가',zorder=3)
    ax.plot(x,sma(p,5),color=RD,lw=1.5,label='5일선',zorder=4)
    ax.plot(x,sma(p,20),color=AC,lw=2.2,label='20일선',zorder=5)
    ax.plot(x,sma(p,60),color=BL,lw=2.2,label='60일선',zorder=5)
    leg(ax); plt.tight_layout(pad=0.4); return save(fig)

def ch_cross():
    np.random.seed(3)
    p=np.concatenate([np.linspace(125,88,70)+np.random.normal(0,1.5,70),
                      np.linspace(88,128,80)+np.random.normal(0,1.5,80)])
    m5=sma(p,5); m20=sma(p,20)
    fig,ax=plt.subplots(figsize=(7,4.5)); fig.patch.set_facecolor(BH); ax_s(ax)
    x=np.arange(len(p))
    ax.plot(x,p,  color=WH,lw=1.0,alpha=0.4,zorder=2)
    ax.plot(x,m5, color=RD,lw=2.0,label='단기선(5일)',zorder=4)
    ax.plot(x,m20,color=BL,lw=2.0,label='장기선(20일)',zorder=4)
    diff=m5-m20; dc=gc=None
    for i in range(22,len(p)-1):
        if np.isnan(diff[i]) or np.isnan(diff[i-1]): continue
        if diff[i-1]>0 and diff[i]<0 and dc is None: dc=i
        if diff[i-1]<0 and diff[i]>0 and gc is None and dc: gc=i
    if dc:
        y=(m5[dc]+m20[dc])/2
        ax.scatter(dc,y,s=180,color=RD,zorder=7,marker='v')
        ax.annotate('데드크로스\n(매도 신호)',xy=(dc,y),xytext=(dc+10,y+7),
                    color=RD,fontsize=9.5,fontweight='bold',
                    arrowprops=dict(arrowstyle='->',color=RD,lw=1.3))
        ax.axvspan(dc,gc or dc+30,alpha=0.07,color=RD)
    if gc:
        y=(m5[gc]+m20[gc])/2
        ax.scatter(gc,y,s=180,color=AC,zorder=7,marker='^')
        ax.annotate('골든크로스\n(매수 신호)',xy=(gc,y),xytext=(gc+10,y-9),
                    color=AC,fontsize=9.5,fontweight='bold',
                    arrowprops=dict(arrowstyle='->',color=AC,lw=1.3))
        ax.axvspan(gc,len(p)-1,alpha=0.07,color=AC)
    leg(ax); plt.tight_layout(pad=0.4); return save(fig)

def ch_align():
    fig,(ax1,ax2)=plt.subplots(1,2,figsize=(10,4.2)); fig.patch.set_facecolor(BH)
    for ax,seed,sign,title,tc in [
        (ax1,1,1,'정배열 (상승 추세)',AC),
        (ax2,2,-1,'역배열 (하락 추세)',RD)]:
        p=trend(100,100,sign*0.25,1.5,seed)
        ax_s(ax,title,tc); x=np.arange(len(p))
        ax.plot(x,sma(p,60),color=BL,lw=2.0,label='60일선')
        ax.plot(x,sma(p,20),color=AC,lw=2.0,label='20일선')
        ax.plot(x,sma(p,5), color=RD,lw=1.8,label='5일선')
        ax.plot(x,p,        color=WH,lw=1.5,label='주가',alpha=0.9,zorder=6)
        leg(ax,loc='upper left' if sign>0 else 'upper right')
    plt.tight_layout(pad=0.4); fig.subplots_adjust(wspace=0.12); return save(fig)

def ch_bb():
    p=ranging(120,100,10,42); m=sma(p,20); s=std_d(p,20)
    ub=m+2*s; lb=m-2*s
    fig,ax=plt.subplots(figsize=(7,4.5)); fig.patch.set_facecolor(BH); ax_s(ax)
    x=np.arange(len(p))
    ax.fill_between(x,lb,ub,alpha=0.12,color=BL)
    ax.plot(x,ub,color=RD,lw=1.8,ls='--',label='상단밴드 (+2σ)')
    ax.plot(x,m, color=AC,lw=2.2,label='중간밴드 (MA20)')
    ax.plot(x,lb,color=BL,lw=1.8,ls='--',label='하단밴드 (-2σ)')
    ax.plot(x,p, color=WH,lw=1.2,alpha=0.85,label='주가',zorder=5)
    # overbought/oversold markers
    for i in range(20,len(p)):
        if np.isnan(ub[i]): continue
        if p[i]>=ub[i]*0.997:
            ax.scatter(i,p[i],s=100,color=RD,zorder=8,marker='v',edgecolors=BH,linewidths=0.5)
        elif p[i]<=lb[i]*1.003:
            ax.scatter(i,p[i],s=100,color=BL,zorder=8,marker='^',edgecolors=BH,linewidths=0.5)
    handles=[
        Line2D([0],[0],color=WH,lw=1.5,label='주가'),
        Line2D([0],[0],color=RD,lw=1.8,ls='--',label='상단밴드(+2σ)'),
        Line2D([0],[0],color=AC,lw=2.2,label='중간밴드(MA20)'),
        Line2D([0],[0],color=BL,lw=1.8,ls='--',label='하단밴드(-2σ)'),
        Line2D([0],[0],color=RD,lw=0,marker='v',ms=8,label='과매수'),
        Line2D([0],[0],color=BL,lw=0,marker='^',ms=8,label='과매도'),
    ]
    leg(ax,handles=handles); plt.tight_layout(pad=0.4); return save(fig)

def ch_bb_dual():
    np.random.seed(10)
    r1=np.concatenate([np.random.normal(0,0.35,65),np.random.normal(0.3,1.6,55)])
    p1=100*np.exp(np.cumsum(r1)/100)
    p2=100*np.exp(np.cumsum(np.random.normal(0.22,1.6,120))/100)
    fig,(ax1,ax2)=plt.subplots(1,2,figsize=(10,4.2)); fig.patch.set_facecolor(BH)
    for ax,p,title,tc in [
        (ax1,p1,'Squeeze -> 큰 움직임 예고',AC),
        (ax2,p2,'Expansion -> 추세 형성 신호',RD)]:
        ax_s(ax,title,tc); x=np.arange(len(p))
        m=sma(p,20); s=std_d(p,20); ub=m+2*s; lb=m-2*s
        ax.fill_between(x,lb,ub,alpha=0.15,color=BL)
        ax.plot(x,ub,color=RD,lw=1.5,ls='--')
        ax.plot(x,m, color=AC,lw=2.0)
        ax.plot(x,lb,color=BL,lw=1.5,ls='--')
        ax.plot(x,p, color=WH,lw=1.5,alpha=0.9,zorder=5)
        if ax==ax1:
            ax.axvline(x=65,color=AC,lw=1.5,ls=':',alpha=0.9)
            ylim=ax.get_ylim()
            ax.text(30,ylim[0]+(ylim[1]-ylim[0])*0.95,'← 수축 구간',color=GY,fontsize=9,va='top')
            ax.text(75,ylim[0]+(ylim[1]-ylim[0])*0.95,'폭발 →',color=AC,fontsize=9,va='top',fontweight='bold')
    plt.tight_layout(pad=0.4); fig.subplots_adjust(wspace=0.12); return save(fig)

def ch_ichi():
    np.random.seed(5); n=180
    p=100*np.exp(np.cumsum(np.random.normal(0.12,1.2,n))/100)
    h=p*(1+np.abs(np.random.normal(0,0.004,n)))
    l=p*(1-np.abs(np.random.normal(0,0.004,n)))
    def hi(a,per,i): return np.max(a[max(0,i-per+1):i+1])
    def lo(a,per,i): return np.min(a[max(0,i-per+1):i+1])
    tk=np.array([(hi(h,9,i)+lo(l,9,i))/2 for i in range(n)])
    kj=np.array([(hi(h,26,i)+lo(l,26,i))/2 for i in range(n)])
    s1r=(tk+kj)/2; s2r=np.array([(hi(h,52,i)+lo(l,52,i))/2 for i in range(n)])
    s1=np.full(n,np.nan); s1[26:]=s1r[:-26]
    s2=np.full(n,np.nan); s2[26:]=s2r[:-26]
    lag=np.full(n,np.nan); lag[:-26]=p[26:]
    fig,ax=plt.subplots(figsize=(8.5,4.8)); fig.patch.set_facecolor(BH); ax_s(ax)
    x=np.arange(n)
    ax.fill_between(x,s1,s2,where=(s1>=s2),alpha=0.28,color=GN)
    ax.fill_between(x,s1,s2,where=(s2>s1), alpha=0.28,color=RD)
    ax.plot(x,p,  color=WH,lw=2.0,label='주가',zorder=6,alpha=0.9)
    ax.plot(x,tk, color=RD,lw=1.5,label='전환선(9일)')
    ax.plot(x,kj, color=BL,lw=2.0,label='기준선(26일)')
    ax.plot(x,s1, color=AC,lw=1.5,label='선행스팬1',ls='--')
    ax.plot(x,s2, color=GN,lw=1.5,label='선행스팬2',ls='--')
    ax.plot(x,lag,color=OR,lw=1.5,label='후행스팬',ls=':')
    handles=[
        Line2D([0],[0],color=WH,lw=2,label='주가'),
        Line2D([0],[0],color=RD,lw=1.5,label='전환선(9)'),
        Line2D([0],[0],color=BL,lw=2,label='기준선(26)'),
        Line2D([0],[0],color=AC,lw=1.5,ls='--',label='선행스팬1'),
        Line2D([0],[0],color=GN,lw=1.5,ls='--',label='선행스팬2'),
        Line2D([0],[0],color=OR,lw=1.5,ls=':',label='후행스팬'),
        Patch(facecolor=GN,alpha=0.4,label='양운(지지)'),
        Patch(facecolor=RD,alpha=0.4,label='음운(저항)'),
    ]
    leg(ax,handles=handles); plt.tight_layout(pad=0.4); return save(fig)

def ch_ichi_cloud():
    np.random.seed(7); n=160
    r=np.concatenate([np.random.normal(-0.1,1.0,55),
                      np.random.normal(0.05,0.7,25),
                      np.random.normal(0.28,1.0,80)])
    p=108*np.exp(np.cumsum(r)/100)
    h=p*1.005; l=p*0.995
    def hi(a,per,i): return np.max(a[max(0,i-per+1):i+1])
    def lo(a,per,i): return np.min(a[max(0,i-per+1):i+1])
    tk=np.array([(hi(h,9,i)+lo(l,9,i))/2 for i in range(n)])
    kj=np.array([(hi(h,26,i)+lo(l,26,i))/2 for i in range(n)])
    s1r=(tk+kj)/2; s2r=np.array([(hi(h,52,i)+lo(l,52,i))/2 for i in range(n)])
    s1=np.full(n,np.nan); s1[26:]=s1r[:-26]
    s2=np.full(n,np.nan); s2[26:]=s2r[:-26]
    fig,ax=plt.subplots(figsize=(7,4.5)); fig.patch.set_facecolor(BH); ax_s(ax)
    x=np.arange(n)
    ax.fill_between(x,s1,s2,where=(s1>=s2),alpha=0.35,color=GN)
    ax.fill_between(x,s1,s2,where=(s2>s1), alpha=0.35,color=RD)
    ax.plot(x,s1,color=AC,lw=1.0,ls='--',alpha=0.6)
    ax.plot(x,s2,color=GN,lw=1.0,ls='--',alpha=0.6)
    ax.plot(x,p, color=WH,lw=2.2,zorder=6)
    # annotations
    ax.annotate('음운(저항)\n하락 추세',xy=(28,p[28]),xytext=(10,p[28]-7),
                color=RD,fontsize=9.5,fontweight='bold',ha='center',
                arrowprops=dict(arrowstyle='->',color=RD,lw=1.3))
    # find breakthrough
    ci=None
    for i in range(85,n-5):
        if np.isnan(s1[i]) or np.isnan(s2[i]): continue
        ct=max(s1[i],s2[i]); ct_prev=max(s1[i-1] if not np.isnan(s1[i-1]) else 0,
                                           s2[i-1] if not np.isnan(s2[i-1]) else 0)
        if p[i]>ct and p[i-1]<=ct_prev and ci is None: ci=i
    if ci:
        ax.scatter(ci,p[ci],s=250,color=AC,zorder=8,marker='^',edgecolors=BH)
        ax.annotate('구름 돌파!\n추세 전환',xy=(ci,p[ci]),xytext=(ci+12,p[ci]+6),
                    color=AC,fontsize=10,fontweight='bold',
                    arrowprops=dict(arrowstyle='->',color=AC,lw=1.5))
    mu=int(n*0.82)
    ax.annotate('양운(지지)\n상승 추세',xy=(mu,p[mu]),xytext=(mu-15,p[mu]-9),
                color=GN,fontsize=9.5,fontweight='bold',ha='center',
                arrowprops=dict(arrowstyle='->',color=GN,lw=1.3))
    handles=[Line2D([0],[0],color=WH,lw=2,label='주가'),
             Patch(facecolor=GN,alpha=0.45,label='양운(지지)'),
             Patch(facecolor=RD,alpha=0.45,label='음운(저항)')]
    leg(ax,handles=handles); plt.tight_layout(pad=0.4); return save(fig)

def ch_env():
    p=ranging(120,100,9,99); m=sma(p,20); pct=0.08
    ub=m*(1+pct); lb=m*(1-pct)
    fig,ax=plt.subplots(figsize=(7,4.5)); fig.patch.set_facecolor(BH); ax_s(ax)
    x=np.arange(len(p))
    ax.fill_between(x,lb,ub,alpha=0.1,color=BL)
    ax.plot(x,ub,color=RD,lw=1.8,ls='--',label='상단선(+8%)')
    ax.plot(x,m, color=AC,lw=2.2,label='중간선(MA20)')
    ax.plot(x,lb,color=BL,lw=1.8,ls='--',label='하단선(-8%)')
    ax.plot(x,p, color=WH,lw=1.5,alpha=0.9,label='주가',zorder=5)
    for i in range(20,len(p)):
        if np.isnan(ub[i]): continue
        if p[i]>=ub[i]*0.997: ax.scatter(i,p[i],s=90,color=RD,zorder=8,marker='v',edgecolors=BH)
        elif p[i]<=lb[i]*1.003: ax.scatter(i,p[i],s=90,color=BL,zorder=8,marker='^',edgecolors=BH)
    handles=[
        Line2D([0],[0],color=WH,lw=1.5,label='주가'),
        Line2D([0],[0],color=RD,lw=1.8,ls='--',label='상단선(+8%)'),
        Line2D([0],[0],color=AC,lw=2.2,label='중간선(MA20)'),
        Line2D([0],[0],color=BL,lw=1.8,ls='--',label='하단선(-8%)'),
        Line2D([0],[0],color=RD,lw=0,marker='v',ms=9,label='과매수 터치'),
        Line2D([0],[0],color=BL,lw=0,marker='^',ms=9,label='과매도 터치'),
    ]
    leg(ax,handles=handles); plt.tight_layout(pad=0.4); return save(fig)

# ═══════════════════════════════════════════════════════
# PPT HELPERS
# ═══════════════════════════════════════════════════════

prs=Presentation(); prs.slide_width=W; prs.slide_height=H

def ns():
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    f=sl.background.fill; f.solid(); f.fore_color.rgb=pBG; return sl

def tx(sl,t,x,y,w,h,sz=16,b=False,c=None,al=PP_ALIGN.LEFT,it=False):
    if c is None: c=pWH
    tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=al; r=p.add_run(); r.text=t
    r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it
    r.font.color.rgb=c; r.font.name='Malgun Gothic'

def txl(sl,lines,x,y,w,h,sz=15):
    tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,(t,c,b) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(3); r=p.add_run(); r.text=t
        r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c
        r.font.name='Malgun Gothic'

def rc(sl,x,y,w,h,fc=None,lc=None,lw=Pt(1)):
    s=sl.shapes.add_shape(1,x,y,w,h)
    if fc: s.fill.solid(); s.fill.fore_color.rgb=fc
    else: s.fill.background()
    if lc: s.line.color.rgb=lc; s.line.width=lw
    else: s.line.fill.background()

def hdr(sl,title,sub=None):
    rc(sl,0,0,W,Inches(0.07),fc=pAC)
    tx(sl,title,Inches(0.5),Inches(0.13),Inches(11),Inches(0.65),sz=31,b=True)
    if sub: tx(sl,sub,Inches(0.5),Inches(0.72),Inches(11.5),Inches(0.38),sz=13,c=pGY)

def pic(sl,buf,x,y,w,h): sl.shapes.add_picture(buf,x,y,w,h)

LX=Inches(0.3); LW=Inches(5.75); RX=Inches(6.3); RW=Inches(6.8); CY=Inches(1.2)

# ═══════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════

# SL01 표지
sl=ns()
rc(sl,0,0,W,Inches(0.1),fc=pAC)
rc(sl,0,H-Inches(0.1),W,Inches(0.1),fc=pAC)
rc(sl,0,0,Inches(0.15),H,fc=pAC)
tx(sl,"차트 지표 기초",Inches(0.5),Inches(1.3),Inches(9),Inches(1.5),sz=54,b=True)
tx(sl,"이동평균선 · 볼린저밴드 · 일목균형표 · 엔벨로프선",
   Inches(0.5),Inches(3.0),Inches(10),Inches(0.7),sz=21,c=pAC)
tx(sl,"주요 기술적 지표 4가지를 한 번에 이해하는 차트 분석 입문 자료입니다.",
   Inches(0.5),Inches(3.85),Inches(9),Inches(0.6),sz=15,c=pGY)
# decorative right chart
buf=ch_ma()
pic(sl,buf,Inches(8.3),Inches(1.6),Inches(4.8),Inches(5.6))

# SL02 목차
sl=ns()
rc(sl,0,0,W,Inches(0.07),fc=pAC)
tx(sl,"목차",Inches(0.5),Inches(0.18),Inches(5),Inches(0.65),sz=34,b=True)
items=[
    ("01","이동평균선","Moving Average","일정 기간 종가 평균선\n추세·지지·저항 파악의 기본",pRD),
    ("02","볼린저밴드","Bollinger Bands","MA ± 2σ 3선 구조\n변동성과 가격 위치 동시 파악",pBL),
    ("03","일목균형표","Ichimoku Cloud","5선+구름대 종합 지표\n추세·지지·저항·모멘텀 한눈에",pAC),
    ("04","엔벨로프선","Envelope","MA ± 고정% 3선 구조\n횡보장 과매수·과매도 파악",pOR),
]
for i,(num,kor,eng,desc,c) in enumerate(items):
    col=i%2; row=i//2
    bx=Inches(0.4)+col*Inches(6.4); by=Inches(1.2)+row*Inches(2.9)
    rc(sl,bx,by,Inches(6.0),Inches(2.65),fc=pD2,lc=c,lw=Pt(0.8))
    rc(sl,bx,by,Inches(0.12),Inches(2.65),fc=c)
    tx(sl,num,bx+Inches(0.22),by+Inches(0.15),Inches(0.8),Inches(0.7),sz=36,b=True,c=c)
    tx(sl,kor,bx+Inches(0.95),by+Inches(0.2),Inches(3),Inches(0.5),sz=22,b=True)
    tx(sl,eng,bx+Inches(0.95),by+Inches(0.65),Inches(3.5),Inches(0.35),sz=13,c=pGY,it=True)
    tx(sl,desc,bx+Inches(0.22),by+Inches(1.1),Inches(5.5),Inches(1.3),sz=14,
       c=RGBColor(0xCC,0xCC,0xCC))

# SL03 이동평균선 개념
sl=ns()
hdr(sl,"이동평균선 (Moving Average)","일정 기간 종가의 평균을 이어 추세 방향을 파악하는 지표")
txl(sl,[
    ("📌 개념",pAC,True),("",pWH,False),
    ("  · 일정 기간(n일) 동안 종가 평균값을 연결한 선",pWH,False),
    ("  · 단기 노이즈를 제거하고 추세 방향을 부드럽게 표시",pWH,False),
    ("  · 기간이 길수록 느리게, 짧을수록 빠르게 반응",pWH,False),
],LX,CY,LW,Inches(2.2))
txl(sl,[
    ("📊 주요 기간선",pAC,True),("",pWH,False),
    ("    5일선  ── 1주일 평균  (초단기)",pWH,False),
    ("   20일선  ── 1달 평균   (가장 많이 쓰이는 기준)",pWH,False),
    ("   60일선  ── 3달 평균   (중기)",pWH,False),
    ("  120일선  ── 6달 평균   (장기)",pWH,False),
    ("  240일선  ── 1년 평균   (초장기)",pWH,False),
],LX,Inches(3.55),LW,Inches(3.4))
pic(sl,ch_ma(),RX,CY,RW,Inches(5.9))

# SL04 크로스 신호
sl=ns()
hdr(sl,"이동평균선 — 크로스 신호","단기선이 장기선을 돌파하는 순간 매매 신호 발생")
txl(sl,[
    ("🟡 골든크로스 (매수 신호)",pAC,True),("",pWH,False),
    ("  단기선이 장기선을 아래 → 위로 돌파",pWH,False),
    ("  상승 모멘텀이 장기 평균을 넘어섰다는 신호",pWH,False),
    ("  상승 추세 전환 알림",pWH,False),
],LX,CY,LW,Inches(2.2))
txl(sl,[
    ("🔴 데드크로스 (매도 신호)",pRD,True),("",pWH,False),
    ("  단기선이 장기선을 위 → 아래로 돌파",pWH,False),
    ("  하락 모멘텀이 장기 평균 밑으로 내려간 신호",pWH,False),
    ("  하락 추세 전환 알림",pWH,False),
],LX,Inches(3.6),LW,Inches(2.2))
txl(sl,[
    ("⚠  주의사항",pOR,True),
    ("  · 후행성 지표 — 신호는 항상 약간 늦게 나온다",pWH,False),
    ("  · 횡보장에서는 크로스가 잦아 신뢰도 하락",pWH,False),
    ("  · 단독보다 다른 지표와 병행 권장",pWH,False),
],LX,Inches(5.9),LW,Inches(1.5))
pic(sl,ch_cross(),RX,CY,RW,Inches(5.9))

# SL05 정배열/역배열 — 풀폭 차트
sl=ns()
hdr(sl,"이동평균선 — 정배열과 역배열","이동평균선 순서로 시장 전체 분위기를 파악한다")
txl(sl,[
    ("  📈 정배열: 주가 > 5일 > 20일 > 60일  →  강한 상승 추세  →  눌림목 매수 유효",pAC,True),
    ("  📉 역배열: 60일 > 20일 > 5일 > 주가  →  강한 하락 추세  →  반등 매도 / 단기 주의",pRD,True),
],LX,Inches(1.05),Inches(12.7),Inches(0.85),sz=14)
pic(sl,ch_align(),Inches(0.3),Inches(1.95),Inches(12.7),Inches(5.3))

# SL06 볼린저밴드 개념 + 구조
sl=ns()
hdr(sl,"볼린저밴드 (Bollinger Bands)","이동평균선 ± 표준편차(2σ)로 만든 3개의 밴드")
txl(sl,[
    ("📌 개념",pAC,True),("",pWH,False),
    ("  존 볼린저(John Bollinger)가 1980년대 개발",pWH,False),
    ("  이동평균 기준으로 통계적 변동 범위를 시각화",pWH,False),
    ("  가격의 변동성과 현재 위치를 동시에 파악",pWH,False),
    ("  통계적으로 가격의 약 95%가 밴드 안에 존재",pGY,False),
],LX,CY,LW,Inches(2.6))
txl(sl,[
    ("🔢 구성 공식",pAC,True),("",pWH,False),
    ("  중간밴드  =  20일 이동평균 (MA20)",pAC,True),
    ("  상단밴드  =  MA20  +  2 × 표준편차(σ)",pRD,True),
    ("  하단밴드  =  MA20  −  2 × 표준편차(σ)",pBL,True),
    ("",pWH,False),
    ("  ※ 20일, 2σ는 기본값 — 조정 가능",pGY,False),
],LX,Inches(4.0),LW,Inches(3.0))
pic(sl,ch_bb(),RX,CY,RW,Inches(5.9))

# SL07 볼린저밴드 수축/팽창 — 풀폭
sl=ns()
hdr(sl,"볼린저밴드 — 수축(Squeeze)과 팽창(Expansion)","밴드 폭의 변화가 변동성과 방향성을 예고한다")
txl(sl,[
    ("  🗜 수축: 밴드 폭 좁아짐 = 변동성 감소 → 큰 움직임 임박 (방향 미확인)",pAC,True),
    ("  ↔ 팽창: 밴드 폭 넓어짐 = 변동성 증가 → 추세 형성 중 (방향 확인 후 대응)",pRD,True),
],LX,Inches(1.05),Inches(12.7),Inches(0.85),sz=14)
pic(sl,ch_bb_dual(),Inches(0.3),Inches(1.95),Inches(12.7),Inches(5.3))

# SL08 일목균형표 5선
sl=ns()
hdr(sl,"일목균형표 (Ichimoku Cloud)","5개 선이 서로 균형을 이루며 추세·지지·저항·모멘텀을 동시 표시")
lines5=[
    ("전환선","(9일 최고+최저) ÷ 2",pRD),
    ("기준선","(26일 최고+최저) ÷ 2",pBL),
    ("선행스팬1","(전환선+기준선) ÷ 2 → 26일 후",pAC),
    ("선행스팬2","(52일 최고+최저) ÷ 2 → 26일 후",pGN),
    ("후행스팬","당일 종가를 26일 전에 표시",pOR),
]
by=Inches(1.2)
for name,formula,c in lines5:
    rc(sl,LX,by,Inches(1.55),Inches(0.65),fc=c)
    tx(sl,name,LX+Inches(0.05),by+Inches(0.1),Inches(1.45),Inches(0.42),
       sz=13,b=True,c=pBG,al=PP_ALIGN.CENTER)
    tx(sl,formula,Inches(2.0),by+Inches(0.1),Inches(4.1),Inches(0.42),sz=13)
    by+=Inches(0.73)
txl(sl,[
    ("💡 구름대",pAC,True),
    ("  선행스팬1 > 2 = 양운(지지)",pGN,False),
    ("  선행스팬2 > 1 = 음운(저항)",pRD,False),
    ("  가격이 구름 위 = 상승 추세",pWH,False),
    ("  가격이 구름 아래 = 하락 추세",pWH,False),
    ("  구름 돌파 = 강한 추세 전환",pAC,True),
],LX,Inches(5.0),LW,Inches(2.3),sz=14)
pic(sl,ch_ichi(),RX,CY,RW,Inches(5.9))

# SL09 일목균형표 구름대 해석
sl=ns()
hdr(sl,"일목균형표 — 구름대(雲) 해석","선행스팬1과 2 사이 구역 = 미래의 지지·저항대")
txl(sl,[
    ("☁ 구름대란?",pAC,True),("",pWH,False),
    ("  선행스팬1과 2 사이 영역",pWH,False),
    ("  현재보다 26일 앞에 그려짐 → 미래 지지·저항 예고",pWH,False),
],LX,CY,LW,Inches(1.8))
txl(sl,[
    ("🟢 양운 (상승 구름)",pGN,True),
    ("  스팬1 > 스팬2 / 가격이 구름 위",pWH,False),
    ("  → 상승 추세 / 강한 지지",pWH,False),
],LX,Inches(3.2),LW,Inches(1.7),sz=14)
txl(sl,[
    ("🔴 음운 (하락 구름)",pRD,True),
    ("  스팬2 > 스팬1 / 가격이 구름 아래",pWH,False),
    ("  → 하락 추세 / 강한 저항",pWH,False),
],LX,Inches(5.0),LW,Inches(1.7),sz=14)
tx(sl,"💡 구름 두꺼울수록 지지·저항 강함 / 구름 돌파 = 추세 전환 신호",
   LX,Inches(6.8),LW,Inches(0.45),sz=13,b=True,c=pAC)
pic(sl,ch_ichi_cloud(),RX,CY,RW,Inches(5.9))

# SL10 엔벨로프선
sl=ns()
hdr(sl,"엔벨로프선 (Envelope)","이동평균에서 고정 비율(%)로 평행 이동 — 횡보장 과매수·과매도 파악")
txl(sl,[
    ("📌 개념 및 구성",pAC,True),("",pWH,False),
    ("  상단선 = MA × (1 + n%)",pRD,True),
    ("  중간선 = MA (이동평균)",pAC,True),
    ("  하단선 = MA × (1 − n%)",pBL,True),
    ("  n%는 종목 변동성에 맞게 설정 (보통 5~15%)",pGY,False),
],LX,CY,LW,Inches(2.7))
txl(sl,[
    ("📊 볼린저밴드와의 차이",pOR,True),("",pWH,False),
    ("  엔벨로프   → 고정 비율 / 밴드 폭 일정",pWH,False),
    ("  볼린저밴드 → 표준편차  / 밴드 폭 변동",pWH,False),
],LX,Inches(4.1),LW,Inches(1.9),sz=14)
txl(sl,[
    ("💡 활용",pAC,True),
    ("  상단 터치 → 과매수 → 차익 매도 검토",pWH,False),
    ("  하단 터치 → 과매도 → 분할 매수 검토",pWH,False),
    ("  횡보(박스권)장에서 특히 유효",pWH,False),
],LX,Inches(6.1),LW,Inches(1.8),sz=14)
pic(sl,ch_env(),RX,CY,RW,Inches(5.9))

# SL11 4가지 비교 정리
sl=ns()
rc(sl,0,0,W,Inches(0.07),fc=pAC)
tx(sl,"4가지 지표 비교 정리",Inches(0.5),Inches(0.18),Inches(10),Inches(0.65),sz=32,b=True)
cols_h=["지표","구성 방식","주요 용도","특이점"]
col_x=[Inches(0.3),Inches(2.75),Inches(5.6),Inches(9.5)]
col_w=[Inches(2.4),Inches(2.8),Inches(3.8),Inches(3.2)]
rc(sl,Inches(0.3),Inches(1.1),W-Inches(0.6),Inches(0.5),fc=pAC)
for h,cx,cw in zip(cols_h,col_x,col_w):
    tx(sl,h,cx+Inches(0.1),Inches(1.15),cw,Inches(0.38),sz=14,b=True,c=pBG,al=PP_ALIGN.CENTER)
rows_d=[
    ("이동평균선","n일 종가 평균선","추세 방향\n지지·저항 파악\n크로스 신호","후행성 지표\n단순·직관적",pRD),
    ("볼린저밴드","MA ± 2σ\n3선 구조","변동성 분석\n과매수·과매도\n수축·팽창 패턴","변동성 반영\n밴드 폭 변동",pBL),
    ("일목균형표","5선 + 구름대","추세·지지·저항\n모멘텀 동시 파악\n미래 예고",  "가장 복잡\n선행 정보 포함",pAC),
    ("엔벨로프선","MA ± 고정%\n3선 구조","과매수·과매도\n횡보장 매매","고정 비율\n설정 단순",pOR),
]
for ri,(name,comp,use,feat,c) in enumerate(rows_d):
    ry=Inches(1.65)+ri*Inches(1.38)
    rc(sl,Inches(0.3),ry,W-Inches(0.6),Inches(1.33),
       fc=RGBColor(0x1A,0x1A,0x1A) if ri%2==0 else pD2)
    rc(sl,Inches(0.3),ry,Inches(0.1),Inches(1.33),fc=c)
    for d,cx,cw in zip([name,comp,use,feat],col_x,col_w):
        tx(sl,d,cx+Inches(0.12),ry+Inches(0.12),cw-Inches(0.2),Inches(1.1),
           sz=14 if d==name else 13,b=(d==name),c=c if d==name else pWH,
           al=PP_ALIGN.CENTER if d==name else PP_ALIGN.LEFT)
tx(sl,"💡 실전 흐름: 이동평균선으로 추세 확인  →  볼린저밴드·엔벨로프로 진입 타이밍  →  일목균형표로 지지·저항 검증",
   Inches(0.4),Inches(7.08),Inches(12.5),Inches(0.35),sz=12,b=True,c=pAC)

# SL12 핵심 요약
sl=ns()
rc(sl,0,0,W,Inches(0.07),fc=pAC)
tx(sl,"오늘의 핵심 요약",Inches(0.5),Inches(0.18),Inches(8),Inches(0.65),sz=32,b=True)
summ=[
    ("이동평균선","추세의 방향을 확인한다.\n골든·데드크로스로 전환 포착.\n정배열=상승 / 역배열=하락.",pRD),
    ("볼린저밴드","변동성을 눈으로 본다.\n수축 후 팽창 = 큰 움직임 예고.\n상단=과매수 / 하단=과매도.",pBL),
    ("일목균형표","하나로 모든 걸 담는다.\n구름 위=상승 / 아래=하락.\n구름 돌파 = 강한 추세 전환.",pAC),
    ("엔벨로프선","정상 범위를 설정한다.\n상단=과매수 / 하단=과매도.\n횡보장에서 특히 효과적.",pOR),
]
for i,(title,body,c) in enumerate(summ):
    col=i%2; row=i//2
    bx=Inches(0.3)+col*Inches(6.55); by=Inches(1.15)+row*Inches(2.9)
    rc(sl,bx,by,Inches(6.2),Inches(2.7),fc=pD2,lc=c,lw=Pt(1.5))
    rc(sl,bx,by,Inches(0.14),Inches(2.7),fc=c)
    tx(sl,title,bx+Inches(0.3),by+Inches(0.18),Inches(5.6),Inches(0.52),sz=22,b=True,c=c)
    tx(sl,body,bx+Inches(0.3),by+Inches(0.75),Inches(5.7),Inches(1.8),sz=15)

# ── Save ──────────────────────────────────────────────
out=r"C:\haessje\라임아카데미\돈복사_지표기초.pptx"
prs.save(out)
print(f"완료: {out}")
print(f"슬라이드: {len(prs.slides)}장")
