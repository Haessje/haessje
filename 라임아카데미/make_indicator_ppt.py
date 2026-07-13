from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG     = RGBColor(0x0D,0x0D,0x0D)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
ACCENT = RGBColor(0xD4,0xFF,0x00)
RED    = RGBColor(0xFF,0x44,0x44)
BLUE   = RGBColor(0x44,0x88,0xFF)
GRAY   = RGBColor(0x77,0x77,0x77)
DARK   = RGBColor(0x1A,0x1A,0x1A)
DARK2  = RGBColor(0x22,0x22,0x22)
ORANGE = RGBColor(0xFF,0xAA,0x00)
GREEN  = RGBColor(0x00,0xCC,0x66)
PURPLE = RGBColor(0xAA,0x66,0xFF)

W = Inches(13.33)
H = Inches(7.5)

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    f = sl.background.fill; f.solid(); f.fore_color.rgb = BG
    return sl

def tx(sl, text, x, y, w, h, size=16, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(x,y,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = 'Malgun Gothic'
    return tb

def txl(sl, lines, x, y, w, h, size=15):
    """lines = [(text, color, bold)]"""
    tb = sl.shapes.add_textbox(x,y,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    for i,(text,color,bold) in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = 'Malgun Gothic'

def rc(sl, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    s = sl.shapes.add_shape(1,x,y,w,h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def ln(sl, x1, y1, x2, y2, color=WHITE, w=Pt(1.5)):
    c = sl.shapes.add_connector(1,x1,y1,x2,y2)
    c.line.color.rgb = color; c.line.width = w

def pl(sl, pts, color=WHITE, w=Pt(1.5)):
    for i in range(len(pts)-1):
        ln(sl, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1], color, w)

def cp(ox,oy,ow,oh):
    def f(cx,cy):
        return ox+cx/100*ow, oy+(1-cy/100)*oh
    return f

def header(sl, title, sub=None):
    rc(sl,0,0,W,Inches(0.07),fill=ACCENT)
    tx(sl,title,Inches(0.5),Inches(0.13),Inches(10),Inches(0.65),size=32,bold=True)
    if sub:
        tx(sl,sub,Inches(0.5),Inches(0.72),Inches(11),Inches(0.38),size=14,color=GRAY)

def chart_bg(sl, x, y, w, h):
    rc(sl, x, y, w, h, fill=DARK)

# ── Build ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = W; prs.slide_height = H

# ════ SL01 표지 ════════════════════════════════════════════════
sl = new_slide(prs)
rc(sl,0,0,W,Inches(0.1),fill=ACCENT)
rc(sl,0,H-Inches(0.1),W,Inches(0.1),fill=ACCENT)
rc(sl,0,0,Inches(0.15),H,fill=ACCENT)
tx(sl,"차트 지표 기초",Inches(0.5),Inches(1.3),Inches(9),Inches(1.5),size=54,bold=True)
tx(sl,"이동평균선 · 볼린저밴드 · 일목균형표 · 엔벨로프선",
   Inches(0.5),Inches(2.9),Inches(10),Inches(0.7),size=21,color=ACCENT)
tx(sl,"주요 기술적 지표 4가지를 한 번에 이해하는 차트 분석 입문 자료입니다.",
   Inches(0.5),Inches(3.7),Inches(9),Inches(0.6),size=15,color=GRAY)

c = cp(Inches(8.5),Inches(1.8),Inches(4.4),Inches(5.0))
price = [(0,40),(10,38),(22,35),(32,42),(42,38),(52,50),(62,55),(72,48),(82,60),(92,68),(100,64)]
pl(sl,[c(*p) for p in price],WHITE,Pt(2))
ma20  = [(0,41),(20,38),(40,41),(60,50),(80,58),(100,63)]
pl(sl,[c(*p) for p in ma20],ACCENT,Pt(2.5))
ub    = [(0,53),(20,50),(40,53),(60,62),(80,70),(100,75)]
lb    = [(0,29),(20,26),(40,29),(60,38),(80,46),(100,51)]
pl(sl,[c(*p) for p in ub],RED,Pt(1.5))
pl(sl,[c(*p) for p in lb],BLUE,Pt(1.5))

# ════ SL02 목차 ════════════════════════════════════════════════
sl = new_slide(prs)
rc(sl,0,0,W,Inches(0.07),fill=ACCENT)
tx(sl,"목차",Inches(0.5),Inches(0.2),Inches(5),Inches(0.65),size=34,bold=True)

items = [
    ("01","이동평균선","Moving Average","일정 기간 종가 평균선\n추세·지지·저항 파악의 기본",RED),
    ("02","볼린저밴드","Bollinger Bands","MA ± 2σ 3선 구조\n변동성과 가격 위치 동시 파악",BLUE),
    ("03","일목균형표","Ichimoku Cloud","5선+구름대 종합 지표\n추세·지지·저항·모멘텀 한눈에",ACCENT),
    ("04","엔벨로프선","Envelope","MA ± 고정% 3선 구조\n횡보장 과매수·과매도 파악",ORANGE),
]
for i,(num,kor,eng,desc,c) in enumerate(items):
    col=i%2; row=i//2
    bx=Inches(0.4)+col*Inches(6.4)
    by=Inches(1.2)+row*Inches(2.9)
    rc(sl,bx,by,Inches(6.0),Inches(2.65),fill=DARK2,line=c,lw=Pt(0.8))
    rc(sl,bx,by,Inches(0.12),Inches(2.65),fill=c)
    tx(sl,num,bx+Inches(0.22),by+Inches(0.15),Inches(0.8),Inches(0.7),size=36,bold=True,color=c)
    tx(sl,kor,bx+Inches(0.95),by+Inches(0.2),Inches(3),Inches(0.5),size=22,bold=True)
    tx(sl,eng,bx+Inches(0.95),by+Inches(0.65),Inches(3.5),Inches(0.35),size=13,color=GRAY,italic=True)
    tx(sl,desc,bx+Inches(0.22),by+Inches(1.1),Inches(5.5),Inches(1.3),size=14,color=RGBColor(0xCC,0xCC,0xCC))

# ════ SL03 이동평균선 개념 + 종류 ════════════════════════════
sl = new_slide(prs)
header(sl,"이동평균선 (Moving Average)","일정 기간 종가의 평균을 이어 추세 방향을 파악하는 지표")

txl(sl,[
    ("📌 개념",ACCENT,True),("",WHITE,False),
    ("  · 일정 기간(n일) 동안 종가의 평균값을 연결한 선",WHITE,False),
    ("  · 단기 노이즈를 제거하고 추세의 방향을 부드럽게 표시",WHITE,False),
    ("  · 기간이 길수록 느리게 반응, 짧을수록 빠르게 반응",WHITE,False),
],Inches(0.4),Inches(1.2),Inches(6.0),Inches(2.2))

txl(sl,[
    ("📊 주요 기간선",ACCENT,True),("",WHITE,False),
    ("    5일선  ── 1주일 평균 (초단기 / 데이트레이더)",WHITE,False),
    ("   20일선  ── 1달 평균 (가장 많이 쓰이는 기준선)",WHITE,False),
    ("   60일선  ── 3달 평균 (중기 추세 파악)",WHITE,False),
    ("  120일선  ── 6달 평균 (장기 추세)",WHITE,False),
    ("  240일선  ── 1년 평균 (장기 방향성)",WHITE,False),
],Inches(0.4),Inches(3.5),Inches(6.0),Inches(3.5))

c = cp(Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
chart_bg(sl,Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
price= [(0,38),(6,42),(12,36),(18,40),(24,45),(30,41),(36,48),(42,52),(48,47),(54,55),
        (60,58),(66,54),(72,60),(78,64),(84,61),(90,67),(96,71),(100,73)]
ma5  = [(0,39),(12,39),(24,43),(36,44),(48,50),(60,56),(72,60),(84,63),(100,72)]
ma20 = [(0,40),(18,40),(36,43),(54,49),(72,58),(90,66),(100,71)]
ma60 = [(0,41),(24,41),(48,44),(72,54),(100,67)]
pl(sl,[c(*p) for p in price],WHITE,Pt(1.5))
pl(sl,[c(*p) for p in ma5],RED,Pt(1.5))
pl(sl,[c(*p) for p in ma20],ACCENT,Pt(2))
pl(sl,[c(*p) for p in ma60],BLUE,Pt(2))
for i,(t,cl) in enumerate([("주가",WHITE),("5일",RED),("20일",ACCENT),("60일",BLUE)]):
    tx(sl,f"─ {t}",Inches(6.9)+i*Inches(1.45),Inches(1.3),Inches(1.4),Inches(0.3),
       size=12,bold=True,color=cl)

# ════ SL04 크로스 신호 ════════════════════════════════════════
sl = new_slide(prs)
header(sl,"이동평균선 — 크로스 신호","단기선이 장기선을 돌파하는 순간이 매매 신호")

txl(sl,[
    ("🟡 골든크로스 (매수 신호)",ACCENT,True),("",WHITE,False),
    ("  단기선이 장기선을 아래 → 위로 돌파",WHITE,False),
    ("  단기 상승 모멘텀이 장기 평균을 넘어섰다는 뜻",WHITE,False),
    ("  상승 추세 전환을 알리는 신호",WHITE,False),
],Inches(0.4),Inches(1.2),Inches(6.0),Inches(2.2))

txl(sl,[
    ("🔴 데드크로스 (매도 신호)",RED,True),("",WHITE,False),
    ("  단기선이 장기선을 위 → 아래로 돌파",WHITE,False),
    ("  단기 하락 모멘텀이 장기 평균 밑으로 내려간 뜻",WHITE,False),
    ("  하락 추세 전환을 알리는 신호",WHITE,False),
],Inches(0.4),Inches(3.6),Inches(6.0),Inches(2.2))

txl(sl,[
    ("⚠  주의사항",ORANGE,True),
    ("  · 후행성 지표 — 신호는 항상 약간 늦게 나온다",WHITE,False),
    ("  · 횡보장에서는 크로스가 잦아 신뢰도 하락",WHITE,False),
    ("  · 단독 사용보다 다른 지표와 병행 권장",WHITE,False),
],Inches(0.4),Inches(5.9),Inches(6.0),Inches(1.5))

c = cp(Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
chart_bg(sl,Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
short_ma=[(0,70),(12,64),(24,55),(36,45),(48,38),(60,42),(72,52),(84,62),(96,70),(100,73)]
long_ma =[(0,68),(18,63),(36,55),(54,46),(72,48),(90,57),(100,65)]
price2  =[(0,72),(8,65),(18,57),(28,47),(38,38),(48,35),(58,42),(68,53),(78,63),(88,70),(100,74)]
pl(sl,[c(*p) for p in price2],WHITE,Pt(1.5))
pl(sl,[c(*p) for p in short_ma],RED,Pt(2))
pl(sl,[c(*p) for p in long_ma],BLUE,Pt(2))

# Dead cross ~30%
dx,dy = c(33,50)
rc(sl,dx-Inches(0.12),dy-Inches(0.12),Inches(0.24),Inches(0.24),fill=RED)
tx(sl,"데드크로스",dx-Inches(0.7),dy-Inches(0.52),Inches(1.6),Inches(0.38),
   size=12,bold=True,color=RED,align=PP_ALIGN.CENTER)

# Golden cross ~54%
gx,gy = c(56,41)
rc(sl,gx-Inches(0.12),gy-Inches(0.12),Inches(0.24),Inches(0.24),fill=ACCENT)
tx(sl,"골든크로스",gx-Inches(0.7),gy+Inches(0.12),Inches(1.6),Inches(0.38),
   size=12,bold=True,color=ACCENT,align=PP_ALIGN.CENTER)

tx(sl,"─ 단기선",Inches(7.0),Inches(1.28),Inches(1.3),Inches(0.3),size=12,bold=True,color=RED)
tx(sl,"─ 장기선",Inches(8.5),Inches(1.28),Inches(1.3),Inches(0.3),size=12,bold=True,color=BLUE)

# ════ SL05 정배열 / 역배열 ═══════════════════════════════════
sl = new_slide(prs)
header(sl,"이동평균선 — 정배열과 역배열","이동평균선의 순서로 시장 분위기 전체를 파악한다")

rc(sl,Inches(0.3),Inches(1.15),Inches(6.1),Inches(6.0),fill=DARK2)
tx(sl,"📈 정배열 (상승 추세)",Inches(0.5),Inches(1.25),Inches(5.5),Inches(0.48),size=19,bold=True,color=ACCENT)
cl = cp(Inches(0.4),Inches(1.9),Inches(5.8),Inches(3.2))
for pts,col,w in [
    ([(0,28),(20,35),(40,43),(60,52),(80,62),(100,70)],BLUE,Pt(2)),
    ([(0,32),(15,39),(30,47),(50,57),(70,66),(100,73)],ACCENT,Pt(2)),
    ([(0,36),(12,43),(28,51),(46,60),(66,69),(100,76)],RED,Pt(2)),
    ([(0,39),(10,46),(24,55),(42,64),(62,72),(100,79)],WHITE,Pt(2.5)),
]:
    pl(sl,[cl(*p) for p in pts],col,w)
txl(sl,[
    ("  주가 > 5일 > 20일 > 60일",ACCENT,True),
    ("  → 강한 상승 추세 확인",WHITE,False),
    ("  → 눌림목 매수 전략 유효",WHITE,False),
],Inches(0.5),Inches(5.3),Inches(5.5),Inches(1.7))

rc(sl,Inches(6.7),Inches(1.15),Inches(6.3),Inches(6.0),fill=DARK2)
tx(sl,"📉 역배열 (하락 추세)",Inches(6.9),Inches(1.25),Inches(5.5),Inches(0.48),size=19,bold=True,color=RED)
cr = cp(Inches(6.8),Inches(1.9),Inches(5.8),Inches(3.2))
for pts,col,w in [
    ([(0,72),(20,65),(40,57),(60,48),(80,38),(100,30)],BLUE,Pt(2)),
    ([(0,76),(15,69),(30,61),(50,51),(70,41),(100,33)],ACCENT,Pt(2)),
    ([(0,79),(12,72),(28,64),(46,54),(66,44),(100,36)],RED,Pt(2)),
    ([(0,82),(10,75),(24,67),(42,57),(62,46),(100,39)],WHITE,Pt(2.5)),
]:
    pl(sl,[cr(*p) for p in pts],col,w)
txl(sl,[
    ("  60일 > 20일 > 5일 > 주가",RED,True),
    ("  → 강한 하락 추세 확인",WHITE,False),
    ("  → 반등 매도 / 단기 대응 주의",WHITE,False),
],Inches(6.9),Inches(5.3),Inches(5.5),Inches(1.7))

for i,(t,col) in enumerate([("─ 주가",WHITE),("─ 5일",RED),("─ 20일",ACCENT),("─ 60일",BLUE)]):
    tx(sl,t,Inches(3.5)+i*Inches(1.6),Inches(7.1),Inches(1.5),Inches(0.3),size=12,bold=True,color=col)

# ════ SL06 볼린저밴드 개념 + 구조 ════════════════════════════
sl = new_slide(prs)
header(sl,"볼린저밴드 (Bollinger Bands)","이동평균선 ± 표준편차(2σ)로 만든 3개의 밴드")

txl(sl,[
    ("📌 개념",ACCENT,True),("",WHITE,False),
    ("  존 볼린저(John Bollinger)가 1980년대 개발",WHITE,False),
    ("  이동평균선을 기준으로 통계적 변동 범위를 표시",WHITE,False),
    ("  가격의 변동성과 현재 위치를 동시에 파악 가능",WHITE,False),
    ("  통계적으로 가격의 약 95%가 밴드 안에 존재",GRAY,False),
],Inches(0.4),Inches(1.2),Inches(6.0),Inches(2.6))

txl(sl,[
    ("🔢 구성 공식",ACCENT,True),("",WHITE,False),
    ("  중간밴드  =  20일 이동평균선 (MA20)",ACCENT,True),
    ("  상단밴드  =  MA20  +  2 × 표준편차(σ)",RED,True),
    ("  하단밴드  =  MA20  −  2 × 표준편차(σ)",BLUE,True),
    ("",WHITE,False),
    ("  ※ 기간(20일)과 표준편차 배수(2)는 변경 가능",GRAY,False),
],Inches(0.4),Inches(4.0),Inches(6.0),Inches(3.2))

c = cp(Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
chart_bg(sl,Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
mid  =[(0,50),(15,48),(30,52),(45,55),(60,52),(75,56),(90,58),(100,60)]
upb  =[(0,72),(15,70),(30,74),(45,78),(60,74),(75,78),(90,80),(100,82)]
lob  =[(0,28),(15,26),(30,30),(45,32),(60,30),(75,34),(90,36),(100,38)]
price3=[(0,48),(8,60),(16,52),(24,68),(32,72),(40,64),(48,44),(56,38),(64,46),(72,55),(80,52),(88,62),(96,48),(100,57)]
pl(sl,[c(*p) for p in upb],RED,Pt(2))
pl(sl,[c(*p) for p in mid],ACCENT,Pt(2.5))
pl(sl,[c(*p) for p in lob],BLUE,Pt(2))
pl(sl,[c(*p) for p in price3],WHITE,Pt(1.5))

px,py = c(100,82); tx(sl,"상단밴드(+2σ)",px+Inches(0.05),py-Inches(0.18),Inches(1.5),Inches(0.35),size=12,bold=True,color=RED)
px,py = c(100,60); tx(sl,"중간밴드(MA20)",px+Inches(0.05),py-Inches(0.18),Inches(1.5),Inches(0.35),size=12,bold=True,color=ACCENT)
px,py = c(100,38); tx(sl,"하단밴드(-2σ)",px+Inches(0.05),py-Inches(0.18),Inches(1.5),Inches(0.35),size=12,bold=True,color=BLUE)

# ════ SL07 볼린저밴드 수축·팽창 + 활용 ══════════════════════
sl = new_slide(prs)
header(sl,"볼린저밴드 — 수축·팽창과 매매 활용","밴드 폭의 변화로 변동성을 예측하고 신호를 잡는다")

rc(sl,Inches(0.3),Inches(1.15),Inches(6.0),Inches(6.0),fill=DARK2)
tx(sl,"🗜 수축(Squeeze) → 큰 움직임 예고",Inches(0.5),Inches(1.25),Inches(5.5),Inches(0.48),size=17,bold=True,color=ACCENT)
txl(sl,[
    ("  밴드 상·하단이 좁아짐 = 변동성 감소",WHITE,False),
    ("  큰 방향성 움직임이 임박했다는 신호",WHITE,False),
    ("  상승/하락 방향은 돌파 이후 확인 필요",GRAY,False),
],Inches(0.5),Inches(1.8),Inches(5.5),Inches(1.3))
cl2 = cp(Inches(0.5),Inches(3.1),Inches(5.5),Inches(2.5))
for pts,col,w in [
    ([(0,75),(25,70),(50,65),(75,62),(100,60)],RED,Pt(1.5)),
    ([(0,55),(25,53),(50,52),(75,51),(100,50)],ACCENT,Pt(2)),
    ([(0,35),(25,36),(50,39),(75,40),(100,40)],BLUE,Pt(1.5)),
    ([(0,54),(15,58),(30,48),(45,55),(60,51),(75,49),(90,51),(100,50)],WHITE,Pt(1.5)),
]:
    pl(sl,[cl2(*p) for p in pts],col,w)
tx(sl,"← 밴드 폭 점점 좁아짐",Inches(1.5),Inches(5.7),Inches(3),Inches(0.35),size=12,color=GRAY)

rc(sl,Inches(6.7),Inches(1.15),Inches(6.3),Inches(6.0),fill=DARK2)
tx(sl,"↔ 팽창(Expansion) → 추세 형성 신호",Inches(6.9),Inches(1.25),Inches(5.8),Inches(0.48),size=17,bold=True,color=RED)
txl(sl,[
    ("  밴드 상·하단이 넓어짐 = 변동성 증가",WHITE,False),
    ("  추세가 형성되고 있다는 신호",WHITE,False),
    ("  방향 확인 후 추세 추종 전략 유효",GRAY,False),
],Inches(6.9),Inches(1.8),Inches(5.5),Inches(1.3))
cr2 = cp(Inches(6.8),Inches(3.1),Inches(5.8),Inches(2.5))
for pts,col,w in [
    ([(0,58),(25,65),(50,72),(75,80),(100,87)],RED,Pt(1.5)),
    ([(0,50),(25,52),(50,56),(75,60),(100,63)],ACCENT,Pt(2)),
    ([(0,42),(25,39),(50,40),(75,40),(100,39)],BLUE,Pt(1.5)),
    ([(0,50),(15,53),(30,58),(48,63),(65,70),(80,75),(100,80)],WHITE,Pt(1.5)),
]:
    pl(sl,[cr2(*p) for p in pts],col,w)
tx(sl,"밴드 폭 넓어짐 →",Inches(9.5),Inches(5.7),Inches(3),Inches(0.35),size=12,color=GRAY)

# 과매수/과매도 hint
txl(sl,[
    ("💡 상단 터치 = 과매수 / 하단 터치 = 과매도",ACCENT,True),
    ("   단, 강한 추세장에서는 밴드를 타고 계속 이동 가능",WHITE,False),
],Inches(0.4),Inches(7.05),Inches(12.5),Inches(0.38))

# ════ SL08 일목균형표 개념 + 5선 ═════════════════════════════
sl = new_slide(prs)
header(sl,"일목균형표 (Ichimoku Cloud)","5개 선과 구름대로 추세·지지·저항·모멘텀을 동시에 파악")

txl(sl,[
    ("📌 개념",ACCENT,True),("",WHITE,False),
    ("  일본 기자 '일목산인'이 1969년 발표",WHITE,False),
    ("  다섯 개의 선이 서로 균형을 이루는 구조",WHITE,False),
    ("  하나의 지표만으로 추세·지지·저항·모멘텀을 모두 파악",WHITE,False),
    ("  '균형'이 깨지는 지점이 매매 신호",WHITE,False),
],Inches(0.4),Inches(1.2),Inches(6.0),Inches(2.4))

lines5=[
    ("전환선","(9일 최고 + 최저) ÷ 2",RED),
    ("기준선","(26일 최고 + 최저) ÷ 2",BLUE),
    ("선행스팬1","(전환선 + 기준선) ÷ 2  → 26일 후 표시",ACCENT),
    ("선행스팬2","(52일 최고 + 최저) ÷ 2  → 26일 후 표시",GREEN),
    ("후행스팬","당일 종가를 26일 '전'에 표시",ORANGE),
]
by = Inches(3.8)
for i,(name,formula,col) in enumerate(lines5):
    ry = by + i*Inches(0.68)
    rc(sl,Inches(0.4),ry,Inches(1.6),Inches(0.62),fill=col)
    tx(sl,name,Inches(0.45),ry+Inches(0.1),Inches(1.5),Inches(0.42),size=13,bold=True,color=BG,align=PP_ALIGN.CENTER)
    tx(sl,formula,Inches(2.1),ry+Inches(0.1),Inches(4.3),Inches(0.42),size=13,color=WHITE)

c = cp(Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
chart_bg(sl,Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
price4=[(10,35),(22,40),(34,38),(46,45),(58,52),(70,58),(82,63),(92,68),(100,70)]
tenkan=[(10,36),(22,39),(34,39),(46,44),(58,51),(70,57),(82,62),(100,69)]
kijun =[(10,37),(26,38),(42,40),(58,47),(74,54),(90,62),(100,67)]
span1 =[(0,25),(15,28),(30,32),(45,38),(60,44),(75,52),(90,58),(100,62)]
span2 =[(0,20),(20,22),(40,26),(60,32),(80,42),(100,52)]
lagg  =[(20,35),(35,40),(50,38),(65,45),(82,52),(100,58)]
pl(sl,[c(*p) for p in price4],WHITE,Pt(2))
pl(sl,[c(*p) for p in tenkan],RED,Pt(1.5))
pl(sl,[c(*p) for p in kijun],BLUE,Pt(2))
pl(sl,[c(*p) for p in span1],ACCENT,Pt(1.5))
pl(sl,[c(*p) for p in span2],GREEN,Pt(1.5))
pl(sl,[c(*p) for p in lagg],ORANGE,Pt(1.5))
# Cloud fill (simple rect approximation)
for i in range(len(span1)-1):
    x1,y1=c(span1[i][0],span1[i][1])
    x2,y2=c(span1[i+1][0],span1[i+1][1])
    si=min(i,len(span2)-2)
    _,b1=c(span2[si][0],span2[si][1])
    top=min(y1,b1); bot=max(y1,b1)
    if bot>top:
        rc(sl,x1,top,x2-x1,bot-top,fill=RGBColor(0x11,0x33,0x22))

for i,(t,col) in enumerate([("전환",RED),("기준",BLUE),("스팬1",ACCENT),("스팬2",GREEN),("후행",ORANGE)]):
    tx(sl,f"─{t}",Inches(6.9)+i*Inches(1.18),Inches(1.28),Inches(1.1),Inches(0.3),size=11,bold=True,color=col)

# ════ SL09 일목균형표 구름대 ══════════════════════════════════
sl = new_slide(prs)
header(sl,"일목균형표 — 구름대(雲) 해석","선행스팬1과 2 사이 구역이 미래의 지지·저항대")

txl(sl,[
    ("☁ 구름대란?",ACCENT,True),("",WHITE,False),
    ("  선행스팬1과 선행스팬2 사이의 영역",WHITE,False),
    ("  현재보다 26일 앞에 그려짐 → 미래 지지/저항 예고",WHITE,False),
],Inches(0.4),Inches(1.2),Inches(6.0),Inches(1.8))

txl(sl,[
    ("🟢 양운 (상승 구름)",GREEN,True),
    ("  선행스팬1 > 선행스팬2",WHITE,False),
    ("  주가가 구름 위에 있음 = 상승 추세",WHITE,False),
    ("  구름이 두꺼울수록 강한 지지대",WHITE,False),
],Inches(0.4),Inches(3.2),Inches(6.0),Inches(1.9))

txl(sl,[
    ("🔴 음운 (하락 구름)",RED,True),
    ("  선행스팬1 < 선행스팬2",WHITE,False),
    ("  주가가 구름 아래 있음 = 하락 추세",WHITE,False),
    ("  구름이 두꺼울수록 강한 저항대",WHITE,False),
],Inches(0.4),Inches(5.2),Inches(6.0),Inches(1.9))

tx(sl,"💡 구름 돌파 = 강한 추세 전환 신호",
   Inches(0.4),Inches(7.05),Inches(6.0),Inches(0.38),size=13,bold=True,color=ACCENT)

c = cp(Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))
chart_bg(sl,Inches(6.8),Inches(1.2),Inches(6.1),Inches(5.9))

# Bearish cloud (left portion, price below)
bs1=[(0,68),(20,65),(40,62)]
bs2=[(0,78),(20,76),(40,74)]
for i in range(len(bs1)-1):
    x1,y1=c(bs1[i][0],bs1[i][1])
    x2,_=c(bs1[i+1][0],bs1[i+1][1])
    _,b1=c(bs2[i][0],bs2[i][1])
    top=min(y1,b1); bot=max(y1,b1)
    if bot>top: rc(sl,x1,top,x2-x1,bot-top,fill=RGBColor(0x44,0x11,0x11))
pl(sl,[c(*p) for p in bs1],RED,Pt(1.5))
pl(sl,[c(*p) for p in bs2],ORANGE,Pt(1.5))

# Price below bearish cloud
pb=[(0,55),(12,52),(24,48),(36,45),(42,48),(50,52)]
pl(sl,[c(*p) for p in pb],WHITE,Pt(2))

# Bullish cloud (right portion, price above)
gs1=[(55,42),(70,47),(85,54),(100,62)]
gs2=[(55,30),(70,34),(85,40),(100,48)]
for i in range(len(gs1)-1):
    x1,y1=c(gs1[i][0],gs1[i][1])
    x2,_=c(gs1[i+1][0],gs1[i+1][1])
    _,b1=c(gs2[i][0],gs2[i][1])
    top=min(y1,b1); bot=max(y1,b1)
    if bot>top: rc(sl,x1,top,x2-x1,bot-top,fill=RGBColor(0x11,0x33,0x11))
pl(sl,[c(*p) for p in gs1],ACCENT,Pt(1.5))
pl(sl,[c(*p) for p in gs2],GREEN,Pt(1.5))

# Price above bullish cloud
pa=[(52,50),(60,56),(70,62),(80,68),(90,73),(100,76)]
pl(sl,[c(*p) for p in pa],WHITE,Pt(2))

# Transition arrow zone
ax,ay=c(50,50)
tx(sl,"⬆ 구름 돌파",ax-Inches(0.4),ay-Inches(0.45),Inches(1.3),Inches(0.38),
   size=13,bold=True,color=ACCENT,align=PP_ALIGN.CENTER)

tx(sl,"음운 (저항)",Inches(7.2),c(20,74)[1]-Inches(0.02),Inches(1.4),Inches(0.35),size=12,bold=True,color=RED)
tx(sl,"양운 (지지)",Inches(11.2),c(80,37)[1]-Inches(0.02),Inches(1.4),Inches(0.35),size=12,bold=True,color=GREEN)

# ════ SL10 엔벨로프선 개념 + 활용 ═══════════════════════════
sl = new_slide(prs)
header(sl,"엔벨로프선 (Envelope)","이동평균선에서 고정 비율(%)로 평행 이동 — 횡보장 과매수·과매도 파악")

txl(sl,[
    ("📌 개념",ACCENT,True),("",WHITE,False),
    ("  이동평균선을 기준으로 위아래로 n%씩 평행 이동",WHITE,False),
    ("  상단선 = MA × (1 + n%)",RED,True),
    ("  중간선 = MA",ACCENT,True),
    ("  하단선 = MA × (1 − n%)",BLUE,True),
    ("  n%는 종목 변동성에 맞게 조정 (보통 5~15%)",GRAY,False),
],Inches(0.4),Inches(1.2),Inches(5.9),Inches(3.2))

txl(sl,[
    ("📊 볼린저밴드와 차이",ORANGE,True),("",WHITE,False),
    ("  엔벨로프   → 고정 비율 적용 / 밴드 폭 일정",WHITE,False),
    ("  볼린저밴드 → 표준편차 적용 / 밴드 폭 변동",WHITE,False),
],Inches(0.4),Inches(4.6),Inches(5.9),Inches(1.8))

txl(sl,[
    ("💡 활용 포인트",ACCENT,True),
    ("  상단 터치 → 과매수 → 차익 매도 검토",WHITE,False),
    ("  하단 터치 → 과매도 → 분할 매수 검토",WHITE,False),
    ("  횡보(박스권)장에서 특히 유효",WHITE,False),
    ("  변동성 큰 종목 10~15% / 작은 종목 3~5%",GRAY,False),
],Inches(0.4),Inches(6.5),Inches(5.9),Inches(0.92))

c = cp(Inches(6.5),Inches(1.2),Inches(6.5),Inches(5.9))
chart_bg(sl,Inches(6.5),Inches(1.2),Inches(6.5),Inches(5.9))

ma_e  =[(0,50),(20,51),(40,50),(60,51),(80,50),(100,51)]
up_e  =[(0,63),(20,64),(40,63),(60,64),(80,63),(100,64)]
lo_e  =[(0,37),(20,38),(40,37),(60,38),(80,37),(100,38)]
price_e=[(0,50),(8,56),(16,63),(24,60),(32,51),(40,44),(48,37),(56,42),(64,50),(72,56),(80,63),(88,60),(96,52),(100,50)]
pl(sl,[c(*p) for p in up_e],RED,Pt(2))
pl(sl,[c(*p) for p in ma_e],ACCENT,Pt(2.5))
pl(sl,[c(*p) for p in lo_e],BLUE,Pt(2))
pl(sl,[c(*p) for p in price_e],WHITE,Pt(2))

# Markers
for tx_c,ty_c,col,label,loy in [(16,63,RED,"과매수",-0.45),(48,37,BLUE,"과매도",0.1),(80,63,RED,"과매수",-0.45)]:
    mx,my=c(tx_c,ty_c)
    rc(sl,mx-Inches(0.1),my-Inches(0.1),Inches(0.2),Inches(0.2),fill=col)
    tx(sl,label,mx-Inches(0.5),my+Inches(loy),Inches(1.2),Inches(0.38),
       size=12,bold=True,color=col,align=PP_ALIGN.CENTER)

tx(sl,"상단(+n%)",Inches(6.7),c(0,64)[1]-Inches(0.2),Inches(1.2),Inches(0.3),size=11,bold=True,color=RED)
tx(sl,"중간(MA)",Inches(6.7),c(0,51)[1]-Inches(0.2),Inches(1.2),Inches(0.3),size=11,bold=True,color=ACCENT)
tx(sl,"하단(-n%)",Inches(6.7),c(0,38)[1]-Inches(0.02),Inches(1.2),Inches(0.3),size=11,bold=True,color=BLUE)

# ════ SL11 4가지 비교 정리 ═══════════════════════════════════
sl = new_slide(prs)
rc(sl,0,0,W,Inches(0.07),fill=ACCENT)
tx(sl,"4가지 지표 비교 정리",Inches(0.5),Inches(0.18),Inches(10),Inches(0.65),size=32,bold=True)

cols_h=["지표","구성 방식","주요 용도","특이점"]
col_x=[Inches(0.3),Inches(2.75),Inches(5.6),Inches(9.5)]
col_w=[Inches(2.4),Inches(2.8),Inches(3.8),Inches(3.2)]
rc(sl,Inches(0.3),Inches(1.1),W-Inches(0.6),Inches(0.5),fill=ACCENT)
for h,cx,cw in zip(cols_h,col_x,col_w):
    tx(sl,h,cx+Inches(0.1),Inches(1.15),cw,Inches(0.38),size=14,bold=True,color=BG,align=PP_ALIGN.CENTER)

rows_d=[
    ("이동평균선","n일 종가 평균선","추세 방향\n지지·저항 파악\n크로스 신호","후행성 지표\n단순·직관적",RED),
    ("볼린저밴드","MA ± 2σ\n3선 구조","변동성 분석\n과매수·과매도\n수축·팽창 패턴","변동성 반영\n밴드 폭 변동",BLUE),
    ("일목균형표","5선 + 구름대","추세·지지·저항\n모멘텀 동시 파악\n미래 지지/저항 예고","가장 복잡\n선행 정보 포함",ACCENT),
    ("엔벨로프선","MA ± 고정%\n3선 구조","과매수·과매도\n횡보장 매매","고정 비율 적용\n설정 단순",ORANGE),
]
for ri,(name,comp,use,feat,col) in enumerate(rows_d):
    ry=Inches(1.65)+ri*Inches(1.38)
    bg_c=DARK if ri%2==0 else DARK2
    rc(sl,Inches(0.3),ry,W-Inches(0.6),Inches(1.33),fill=bg_c)
    rc(sl,Inches(0.3),ry,Inches(0.1),Inches(1.33),fill=col)
    data=[name,comp,use,feat]
    for j,(d,cx,cw) in enumerate(zip(data,col_x,col_w)):
        tx(sl,d,cx+Inches(0.12),ry+Inches(0.12),cw-Inches(0.2),Inches(1.1),
           size=14 if j==0 else 13,bold=(j==0),
           color=col if j==0 else WHITE,
           align=PP_ALIGN.CENTER if j==0 else PP_ALIGN.LEFT)

tx(sl,"💡 실전 흐름: 이동평균선으로 추세 확인 → 볼린저밴드·엔벨로프로 진입 타이밍 → 일목균형표로 지지·저항 검증",
   Inches(0.4),Inches(7.07),Inches(12.5),Inches(0.35),size=13,bold=True,color=ACCENT)

# ════ SL12 핵심 요약 ══════════════════════════════════════════
sl = new_slide(prs)
rc(sl,0,0,W,Inches(0.07),fill=ACCENT)
tx(sl,"오늘의 핵심 요약",Inches(0.5),Inches(0.18),Inches(8),Inches(0.65),size=32,bold=True)

summaries=[
    ("이동평균선","추세의 방향을 확인한다.\n골든·데드크로스로 추세 전환 포착.\n정배열=상승 / 역배열=하락.",RED),
    ("볼린저밴드","변동성을 눈으로 본다.\n수축 후 팽창 = 큰 움직임 예고.\n상단=과매수 / 하단=과매도.",BLUE),
    ("일목균형표","하나로 모든 걸 담는다.\n구름 위=상승 / 아래=하락.\n구름 돌파 = 강한 추세 전환.",ACCENT),
    ("엔벨로프선","정상 범위를 설정한다.\n상단=과매수 / 하단=과매도.\n횡보장에서 특히 효과적.",ORANGE),
]
for i,(title,body,col) in enumerate(summaries):
    c2=i%2; r2=i//2
    bx=Inches(0.3)+c2*Inches(6.55)
    by=Inches(1.15)+r2*Inches(2.9)
    rc(sl,bx,by,Inches(6.2),Inches(2.7),fill=DARK2,line=col,lw=Pt(1.5))
    rc(sl,bx,by,Inches(0.14),Inches(2.7),fill=col)
    tx(sl,title,bx+Inches(0.3),by+Inches(0.18),Inches(5.6),Inches(0.52),size=22,bold=True,color=col)
    tx(sl,body,bx+Inches(0.3),by+Inches(0.75),Inches(5.7),Inches(1.8),size=15,color=WHITE)

# ── Save ──────────────────────────────────────────────────────
out = r"C:\haessje\라임아카데미\돈복사_지표기초.pptx"
prs.save(out)
print(f"완료: {out}")
print(f"슬라이드: {len(prs.slides)}장")
