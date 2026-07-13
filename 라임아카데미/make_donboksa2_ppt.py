import io
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG    = RGBColor(0x0D, 0x0D, 0x0D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIME  = RGBColor(0xD4, 0xFF, 0x00)
GRAY  = RGBColor(0x88, 0x88, 0x88)
RED   = RGBColor(0xFF, 0x44, 0x44)
BLUE  = RGBColor(0x44, 0x88, 0xFF)
GREEN = RGBColor(0x00, 0xCC, 0x66)
YELLOW= RGBColor(0xFF, 0xCC, 0x00)

W = Inches(13.33)
H = Inches(7.5)

plt.rcParams['font.family'] = 'Malgun Gothic'
plt.rcParams['axes.unicode_minus'] = False

def hr(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def txt(slide, text, left, top, width, height,
        size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.name  = 'Malgun Gothic'
    return txb

def box(slide, left, top, width, height, fill):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def add_img(slide, fig, left, top, width, height):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150,
                bbox_inches='tight', facecolor=fig.get_facecolor())
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width, height)
    plt.close(fig)

def apply_table_style(tbl, col_count, header_col_colors=None):
    for (row, col), cell in tbl.get_celld().items():
        cell.set_edgecolor('#333333')
        if row == 0:
            cell.set_facecolor('#1A1A1A')
            cell.set_text_props(color='#D4FF00', fontweight='bold', fontname='Malgun Gothic')
        else:
            bg = '#1A1A1A'
            if header_col_colors and col in header_col_colors:
                bg = header_col_colors[col]
            cell.set_facecolor(bg)
            cell.set_text_props(color='white', fontname='Malgun Gothic')

# ─────────────────────────────────────────────────────────
# Slide 1 — 타이틀
# ─────────────────────────────────────────────────────────
def slide_title(prs):
    sl = blank_slide(prs)
    box(sl, Inches(0), Inches(0), Inches(0.15), H, LIME)
    txt(sl, '2주차', Inches(0.4), Inches(1.4), Inches(4), Inches(0.7),
        size=18, bold=True, color=LIME)
    txt(sl, '좋은 종목을\n어떻게 고르는가?',
        Inches(0.4), Inches(2.0), Inches(9), Inches(2.8),
        size=52, bold=True, color=WHITE)
    txt(sl, '재무 건전성 스코어링 + 단타 타점 입문',
        Inches(0.4), Inches(4.7), Inches(10), Inches(0.6),
        size=22, color=GRAY)
    txt(sl, '돈복사', Inches(10), Inches(6.7), Inches(3), Inches(0.5),
        size=15, color=GRAY, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────
# Slide 2 — 오늘의 흐름
# ─────────────────────────────────────────────────────────
def slide_agenda(prs):
    sl = blank_slide(prs)
    txt(sl, '오늘의 흐름', Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
        size=14, color=GRAY)
    txt(sl, '2시간, 이렇게 씁니다',
        Inches(0.5), Inches(0.85), Inches(11), Inches(0.7),
        size=32, bold=True, color=WHITE)

    rows = [
        ('0~10분',   '오프닝',    '과제 종목 꺼내기',                GRAY,   '#222222'),
        ('10~40분',  '교육 ①',   '좋은 종목 — 재무 건전성 스코어링', LIME,   '#0D1A00'),
        ('40~55분',  '짝활동 ①', '내 종목, 기준 통과하는가?',        YELLOW, '#1A1500'),
        ('55~85분',  '교육 ②',   '단타 기초 — 타점 3가지',           LIME,   '#0D1A00'),
        ('85~100분', '짝활동 ②', '통과 종목에서 타점 직접 찾기',      YELLOW, '#1A1500'),
        ('100~120분','마무리',    '짝별 공유 + 다음 주 예고',          GRAY,   '#222222'),
    ]

    y = Inches(1.85)
    for time, label, desc, col, bg in rows:
        box(sl, Inches(0.5), y, Inches(12.3), Inches(0.65), hr(bg))
        box(sl, Inches(0.5), y, Inches(0.07), Inches(0.65), col)
        txt(sl, time,  Inches(0.7),  y+Inches(0.1), Inches(1.4), Inches(0.5), size=13, color=GRAY)
        txt(sl, label, Inches(2.2),  y+Inches(0.1), Inches(1.8), Inches(0.5), size=15, bold=True, color=col)
        txt(sl, desc,  Inches(4.1),  y+Inches(0.1), Inches(8.5), Inches(0.5), size=15, color=WHITE)
        y += Inches(0.8)

# ─────────────────────────────────────────────────────────
# Slide 3 — 오프닝
# ─────────────────────────────────────────────────────────
def slide_opening(prs):
    sl = blank_slide(prs)
    txt(sl, '오프닝  |  0~10분', Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
        size=14, color=GRAY)
    txt(sl, '과제 꺼내기', Inches(0.5), Inches(0.85), Inches(10), Inches(0.65),
        size=36, bold=True, color=WHITE)
    txt(sl, '일주일간 매매했던 종목 3가지와 산 이유를 꺼내주세요',
        Inches(0.5), Inches(1.7), Inches(12), Inches(0.6),
        size=20, color=LIME)

    fig, ax = plt.subplots(figsize=(11, 3.2))
    fig.patch.set_facecolor('#0D0D0D')
    ax.set_facecolor('#0D0D0D')
    ax.axis('off')

    data = [['종목 A', '', '', ''],
            ['종목 B', '', '', ''],
            ['종목 C', '', '', '']]
    cols = ['종목명', '산 이유 ①', '산 이유 ②', '산 이유 ③']
    tbl = ax.table(cellText=data, colLabels=cols, cellLoc='center', loc='center')
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(13)
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor('#333333')
        cell.set_height(0.22)
        if r == 0:
            cell.set_facecolor('#1A1A1A')
            cell.set_text_props(color='#D4FF00', fontweight='bold', fontname='Malgun Gothic')
        else:
            cell.set_facecolor('#141414')
            cell.set_text_props(color='white', fontname='Malgun Gothic')
    plt.tight_layout(pad=0.3)
    add_img(sl, fig, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.5))

# ─────────────────────────────────────────────────────────
# Slide 4 — 교육① 인트로
# ─────────────────────────────────────────────────────────
def slide_edu1_intro(prs):
    sl = blank_slide(prs)
    txt(sl, '교육 ①  |  10~40분', Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
        size=14, color=GRAY)
    txt(sl, '좋은 종목을 어떻게 구별하나?',
        Inches(0.5), Inches(0.85), Inches(11), Inches(0.7),
        size=32, bold=True, color=WHITE)

    bad = [
        '뉴스에서 좋다고 하니까',
        '커뮤니티에서 추천해줘서',
        '"이 종목 왠지 오를 것 같아서"',
    ]
    y = Inches(2.0)
    for b in bad:
        txt(sl, '✗  ' + b, Inches(0.7), y, Inches(10), Inches(0.6), size=22, color=GRAY)
        y += Inches(0.75)

    box(sl, Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.9), hr('#0D1A00'))
    box(sl, Inches(0.5), Inches(4.1), Inches(0.07), Inches(1.9), LIME)
    txt(sl, '재무 건전성 스코어링 시스템',
        Inches(0.8), Inches(4.25), Inches(10), Inches(0.55),
        size=22, bold=True, color=LIME)
    txt(sl, '7가지 항목을 숫자로 채점 → SAFE / CAUTION / DANGER 등급으로 판정\n감이 아닌 기준으로, 나쁜 종목을 먼저 걸러낸다',
        Inches(0.8), Inches(4.8), Inches(11.8), Inches(0.9),
        size=17, color=WHITE)

# ─────────────────────────────────────────────────────────
# Slide 5 — 7항목 테이블
# ─────────────────────────────────────────────────────────
def slide_scoring_table(prs):
    sl = blank_slide(prs)
    txt(sl, '재무 건전성 7항목 스코어링',
        Inches(0.5), Inches(0.25), Inches(12), Inches(0.65),
        size=30, bold=True, color=WHITE)

    fig, ax = plt.subplots(figsize=(12, 5.2))
    fig.patch.set_facecolor('#0D0D0D')
    ax.set_facecolor('#0D0D0D')
    ax.axis('off')

    data = [
        ['영업이익 연속 흑자', '3년 이상  +3', '1~2년  +1', '적자  -3'],
        ['부채비율',           '100% 미만  +2', '100~200%  0', '200% 초과  -2'],
        ['유동비율',           '200% 이상  +2', '100~200%  0', '100% 미만  -3'],
        ['매출증가율',         '양수  +1', '—', '음수  -1'],
        ['유보율',             '500% 이상  +1', '100~500%  0', '100% 미만  -1'],
        ['영업현금흐름',       '양수  +1', '—', '음수  -1'],
        ['ROE',               '10% 이상  +1', '0~10%  0', '음수  -1'],
    ]
    cols = ['항목', '좋음', '보통', '나쁨']
    tbl = ax.table(cellText=data, colLabels=cols, cellLoc='center', loc='center')
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(13)

    col_bg = {0: '#1A1A1A', 1: '#0D2A0D', 2: '#1A1A1A', 3: '#2A0D0D'}
    col_fc = {0: 'white', 1: '#00CC66', 2: '#888888', 3: '#FF6666'}
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor('#333333')
        cell.set_height(0.115)
        if r == 0:
            cell.set_facecolor('#1A1A1A')
            cell.set_text_props(color='#D4FF00', fontweight='bold', fontname='Malgun Gothic')
        else:
            cell.set_facecolor(col_bg.get(c, '#1A1A1A'))
            cell.set_text_props(color=col_fc.get(c, 'white'), fontname='Malgun Gothic')

    tbl.auto_set_column_width([0,1,2,3])
    plt.tight_layout(pad=0.3)
    add_img(sl, fig, Inches(0.4), Inches(1.05), Inches(12.5), Inches(6.1))

# ─────────────────────────────────────────────────────────
# Slide 6 — 등급 판정
# ─────────────────────────────────────────────────────────
def slide_grade(prs):
    sl = blank_slide(prs)
    txt(sl, '등급 판정', Inches(0.5), Inches(0.3), Inches(10), Inches(0.65),
        size=30, bold=True, color=WHITE)
    txt(sl, '7항목 합산 점수로 3등급 판정',
        Inches(0.5), Inches(0.9), Inches(10), Inches(0.5),
        size=18, color=GRAY)

    grades = [
        ('SAFE',    '+6 이상', '재무적으로 건전. 종목 검토를 계속 진행해도 좋다.',
         '#00CC66', '#0D2A15'),
        ('CAUTION', '0 ~ +5',  '약점 있음. 어떤 항목이 낮은지 직접 확인 후 판단.',
         '#FFCC00', '#2A2200'),
        ('DANGER',  '-1 이하', '재무 리스크 큼. 아무리 차트가 좋아도 신중하게.',
         '#FF4444', '#2A0D0D'),
    ]

    y = Inches(1.6)
    for label, score, desc, col, bg in grades:
        box(sl, Inches(0.5), y, Inches(12.3), Inches(1.4), hr(bg))
        box(sl, Inches(0.5), y, Inches(0.07), Inches(1.4), hr(col))
        txt(sl, label, Inches(0.8), y+Inches(0.1), Inches(2.5), Inches(0.6),
            size=28, bold=True, color=hr(col))
        txt(sl, score, Inches(3.3), y+Inches(0.1), Inches(2), Inches(0.6),
            size=24, bold=True, color=hr(col))
        txt(sl, desc,  Inches(5.5), y+Inches(0.25), Inches(7.1), Inches(0.7),
            size=17, color=WHITE)
        y += Inches(1.6)

    txt(sl, '* 재무 건전성 = "제외 기준"입니다. SAFE라고 무조건 매수 신호가 아닙니다.',
        Inches(0.5), Inches(6.55), Inches(12), Inches(0.6),
        size=14, color=GRAY)

# ─────────────────────────────────────────────────────────
# Slide 7 — 스코어링 예시
# ─────────────────────────────────────────────────────────
def slide_example(prs):
    sl = blank_slide(prs)
    txt(sl, '스코어링 예시 — 가상 기업 A사',
        Inches(0.5), Inches(0.25), Inches(11), Inches(0.65),
        size=30, bold=True, color=WHITE)

    fig, ax = plt.subplots(figsize=(10, 4.6))
    fig.patch.set_facecolor('#0D0D0D')
    ax.set_facecolor('#0D0D0D')
    ax.axis('off')

    ex = [
        ['영업이익 연속 흑자', '3년 연속 흑자', '+3'],
        ['부채비율',           '85%',           '+2'],
        ['유동비율',           '150%',          '0'],
        ['매출증가율',         '+12%',          '+1'],
        ['유보율',             '320%',          '0'],
        ['영업현금흐름',       '양수',          '+1'],
        ['ROE',               '8%',            '0'],
        ['합계',               '',              '+7  →  SAFE'],
    ]

    tbl = ax.table(cellText=ex, colLabels=['항목', '수치', '점수'],
                   cellLoc='center', loc='center')
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(14)

    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor('#333333')
        cell.set_height(0.105)
        if r == 0:
            cell.set_facecolor('#1A1A1A')
            cell.set_text_props(color='#D4FF00', fontweight='bold', fontname='Malgun Gothic')
        elif r == 8:
            cell.set_facecolor('#0D2A15')
            cell.set_text_props(color='#00CC66', fontweight='bold', fontname='Malgun Gothic')
        else:
            cell.set_facecolor('#1A1A1A')
            score = ex[r-1][2]
            if score.startswith('+') and score != '+0':
                fc = '#00CC66'
            elif score == '0':
                fc = '#FFCC00'
            else:
                fc = '#FF6666'
            if c == 2:
                cell.set_text_props(color=fc, fontweight='bold', fontname='Malgun Gothic')
            else:
                cell.set_text_props(color='white', fontname='Malgun Gothic')

    tbl.auto_set_column_width([0,1,2])
    plt.tight_layout(pad=0.3)
    add_img(sl, fig, Inches(1.2), Inches(1.1), Inches(11), Inches(5.8))

# ─────────────────────────────────────────────────────────
# Slide 8 — 짝활동①
# ─────────────────────────────────────────────────────────
def slide_pair1(prs):
    sl = blank_slide(prs)
    box(sl, Inches(0), Inches(0), W, Inches(0.08), YELLOW)
    txt(sl, '짝활동 ①  |  40~55분',
        Inches(0.5), Inches(0.25), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '내 종목, 기준을 통과하는가?',
        Inches(0.5), Inches(0.8), Inches(12), Inches(0.75),
        size=34, bold=True, color=WHITE)

    box(sl, Inches(0.5), Inches(1.75), Inches(12.3), Inches(1.0), hr('#1A1500'))
    txt(sl, '미션  —  과제 종목 3개 중 하나를 골라 짝이랑 함께 채점하기',
        Inches(0.7), Inches(1.85), Inches(11.8), Inches(0.7),
        size=20, bold=True, color=YELLOW)

    steps = [
        '① 각자 종목 1개 선택 (과제 3개 중)',
        '② 네이버 금융 / 증권사 앱에서 재무 수치 확인',
        '③ 7항목 점수 계산 → 합산',
        '④ SAFE / CAUTION / DANGER 등급 판정',
        '⑤ 짝에게 결과 설명하기  ("이 종목은 ___ 이유로 ___ 등급")',
    ]
    y = Inches(2.95)
    for s in steps:
        txt(sl, s, Inches(0.8), y, Inches(11.5), Inches(0.55), size=20, color=WHITE)
        y += Inches(0.65)

    txt(sl, '⏱  15분', Inches(10.2), Inches(6.75), Inches(2.6), Inches(0.5),
        size=22, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────
# Slide 9 — 교육② 단타 개념
# ─────────────────────────────────────────────────────────
def slide_danta_intro(prs):
    sl = blank_slide(prs)
    txt(sl, '교육 ②  |  55~85분',
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '단타 기초 — 타점을 잡는다는 것',
        Inches(0.5), Inches(0.85), Inches(12), Inches(0.65),
        size=32, bold=True, color=WHITE)

    # Left column — 중장기
    box(sl, Inches(0.5), Inches(1.8), Inches(5.7), Inches(4.5), hr('#1A1A1A'))
    txt(sl, '중장기 투자', Inches(0.7), Inches(1.95), Inches(5), Inches(0.55),
        size=20, bold=True, color=GRAY)
    lt = ['보유기간: 수개월 ~ 수년',
          '기업 가치 성장에 베팅',
          '"어떤 기업"이 핵심',
          '타점보다 종목 선정이 중요']
    y = Inches(2.6)
    for s in lt:
        txt(sl, '• ' + s, Inches(0.8), y, Inches(5.2), Inches(0.5), size=16, color=GRAY)
        y += Inches(0.6)

    # Right column — 단타
    box(sl, Inches(6.5), Inches(1.8), Inches(6.3), Inches(4.5), hr('#0D1A00'))
    box(sl, Inches(6.5), Inches(1.8), Inches(0.07), Inches(4.5), LIME)
    txt(sl, '단타 매매', Inches(6.75), Inches(1.95), Inches(5.5), Inches(0.55),
        size=20, bold=True, color=LIME)
    st = ['보유기간: 당일 ~ 수일',
          '가격 움직임의 패턴에 베팅',
          '"언제"가 "무엇"만큼 중요',
          '좋은 종목 + 좋은 타점 = 핵심']
    y = Inches(2.6)
    for s in st:
        txt(sl, '• ' + s, Inches(6.9), y, Inches(5.7), Inches(0.5), size=16, color=WHITE)
        y += Inches(0.6)

    txt(sl, '오늘 목표  →  재무 건전성을 통과한 종목에서 단타 타점 잡기',
        Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
        size=16, bold=True, color=LIME)

# ─────────────────────────────────────────────────────────
# Slide 10 — 타점 3가지 차트
# ─────────────────────────────────────────────────────────
def slide_danta_chart(prs):
    sl = blank_slide(prs)
    txt(sl, '진입 타점 3가지',
        Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
        size=30, bold=True, color=WHITE)

    np.random.seed(42)
    n = 90
    r = np.random.normal(0.001, 0.016, n)
    C = 100 * np.exp(np.cumsum(r))
    O = np.empty(n); O[0] = C[0]
    for i in range(1, n): O[i] = C[i-1]
    wick = np.abs(C) * 0.01
    H = np.maximum(O, C) + np.abs(np.random.normal(0, wick, n))
    L = np.minimum(O, C) - np.abs(np.random.normal(0, wick, n))

    # 거래량 — 급증 구간 강조
    vol = np.abs(np.random.normal(1, 0.35, n)) + 0.4
    vol[19:23] *= 5.5   # ① 거래량 급증
    vol[47:50] *= 3.0   # ② 눌림목
    vol[67:71] *= 4.5   # ③ 돌파

    # 눌림목 연출 — 가격 일부러 낮춤
    C[45:50] -= 3; O[45:50] -= 3; H[45:50] -= 3; L[45:50] -= 3

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(12, 4.6),
                                    gridspec_kw={'height_ratios': [3, 1]},
                                    facecolor='#0D0D0D')
    for ax in [ax1, ax2]:
        ax.set_facecolor('#0D0D0D')
        for s in ax.spines.values(): s.set_color('#333333')
        ax.tick_params(colors='#555555', labelsize=8)

    x = np.arange(n)
    bull = C >= O; bear = ~bull
    RD = '#FF4444'; BL = '#4488FF'

    for mask, c in [(bull, RD), (bear, BL)]:
        if mask.any():
            ht = np.maximum(np.abs(C[mask]-O[mask]), np.abs(C[mask])*0.001)
            ax1.bar(x[mask], ht, bottom=np.minimum(O[mask],C[mask]),
                    width=0.6, color=c, zorder=4, linewidth=0)
            ax1.vlines(x[mask], L[mask], H[mask], color=c, lw=0.7, zorder=3)

    p1, p2, p3 = 21, 48, 69
    ax1.annotate('① 거래량급증+양봉', xy=(p1, H[p1]+0.3),
                 xytext=(p1-12, H[p1]+5),
                 fontsize=9.5, color='#D4FF00', fontfamily='Malgun Gothic',
                 arrowprops=dict(arrowstyle='->', color='#D4FF00', lw=1.3))
    ax1.annotate('② 눌림목 반등', xy=(p2, L[p2]-0.3),
                 xytext=(p2+3, L[p2]-7),
                 fontsize=9.5, color='#FF9900', fontfamily='Malgun Gothic',
                 arrowprops=dict(arrowstyle='->', color='#FF9900', lw=1.3))
    ax1.annotate('③ 전고점 돌파', xy=(p3, H[p3]+0.3),
                 xytext=(p3-14, H[p3]+5),
                 fontsize=9.5, color='#00CC66', fontfamily='Malgun Gothic',
                 arrowprops=dict(arrowstyle='->', color='#00CC66', lw=1.3))
    for xv, c in [(p1,'#D4FF00'),(p2,'#FF9900'),(p3,'#00CC66')]:
        ax1.axvline(xv, color=c, lw=0.9, linestyle='--', alpha=0.5)

    ax1.set_xlim(-1, n+1)
    ax1.set_ylabel('가격', color='#666666', fontsize=9, fontfamily='Malgun Gothic')
    ax1.yaxis.label.set_color('#666666')

    ax2.bar(x, vol, color=[RD if bull[i] else BL for i in range(n)],
            width=0.6, alpha=0.85)
    for xv, c in [(p1,'#D4FF00'),(p2,'#FF9900'),(p3,'#00CC66')]:
        ax2.axvline(xv, color=c, lw=0.9, linestyle='--', alpha=0.5)
    ax2.set_xlim(-1, n+1)
    ax2.set_ylabel('거래량', color='#666666', fontsize=9, fontfamily='Malgun Gothic')

    plt.tight_layout(pad=0.2)
    add_img(sl, fig, Inches(0.3), Inches(0.95), Inches(12.7), Inches(6.3))

# ─────────────────────────────────────────────────────────
# Slide 11 — 타점 정리
# ─────────────────────────────────────────────────────────
def slide_danta_rules(prs):
    sl = blank_slide(prs)
    txt(sl, '타점 3가지 — 핵심 정리',
        Inches(0.5), Inches(0.25), Inches(11), Inches(0.65),
        size=30, bold=True, color=WHITE)

    rules = [
        ('①', '거래량 급증 + 양봉',
         '평소 거래량의 2배 이상 터지면서 양봉 마감\n"누군가 강하게 사고 있다"는 신호',
         '#D4FF00', '#151A00'),
        ('②', '눌림목 반등',
         '오르던 주가가 잠깐 내려앉다가 다시 튀어오르는 지점\n"올라가던 흐름을 재개"하는 순간',
         '#FF9900', '#1A1200'),
        ('③', '전고점 돌파',
         '이전 고점 가격을 거래량 동반해 넘어서는 순간\n"저항선이 지지선으로 전환"되는 구간',
         '#00CC66', '#001A0D'),
    ]

    y = Inches(1.3)
    for num, title, desc, col, bg in rules:
        box(sl, Inches(0.5), y, Inches(12.3), Inches(1.55), hr(bg))
        box(sl, Inches(0.5), y, Inches(0.07), Inches(1.55), hr(col))
        txt(sl, num,   Inches(0.75), y+Inches(0.1), Inches(0.7), Inches(0.65),
            size=26, bold=True, color=hr(col))
        txt(sl, title, Inches(1.5), y+Inches(0.1), Inches(10), Inches(0.55),
            size=22, bold=True, color=hr(col))
        txt(sl, desc,  Inches(1.5), y+Inches(0.65), Inches(10.8), Inches(0.8),
            size=16, color=WHITE)
        y += Inches(1.8)

# ─────────────────────────────────────────────────────────
# Slide 12 — 짝활동②
# ─────────────────────────────────────────────────────────
def slide_pair2(prs):
    sl = blank_slide(prs)
    box(sl, Inches(0), Inches(0), W, Inches(0.08), YELLOW)
    txt(sl, '짝활동 ②  |  85~100분',
        Inches(0.5), Inches(0.25), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '통과 종목에서 타점 직접 찾기',
        Inches(0.5), Inches(0.8), Inches(12), Inches(0.75),
        size=34, bold=True, color=WHITE)

    box(sl, Inches(0.5), Inches(1.75), Inches(12.3), Inches(1.0), hr('#1A1500'))
    txt(sl, '미션  —  짝활동 ①에서 SAFE / CAUTION 판정 받은 종목의 차트 열기',
        Inches(0.7), Inches(1.85), Inches(11.8), Inches(0.7),
        size=20, bold=True, color=YELLOW)

    steps = [
        '① 차트 열기  (네이버 금융 or 증권사 앱)',
        '② 최근 거래량 급증 구간이 있는가?',
        '③ 눌림목 후 반등 패턴이 보이는가?',
        '④ 전고점 돌파를 시도 중인가?',
        '⑤ 발견한 타점을 짝에게 설명하기',
    ]
    y = Inches(2.95)
    for s in steps:
        txt(sl, s, Inches(0.8), y, Inches(11.5), Inches(0.55), size=20, color=WHITE)
        y += Inches(0.65)

    txt(sl, '⏱  15분', Inches(10.2), Inches(6.75), Inches(2.6), Inches(0.5),
        size=22, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────
# Slide 13 — 마무리
# ─────────────────────────────────────────────────────────
def slide_closing(prs):
    sl = blank_slide(prs)
    txt(sl, '마무리  |  100~120분',
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.55),
        size=14, color=GRAY)
    txt(sl, '오늘 우리가 한 것',
        Inches(0.5), Inches(0.85), Inches(11), Inches(0.65),
        size=36, bold=True, color=WHITE)

    summary = [
        ('재무 건전성 7항목으로', '내 종목을 직접 채점했다'),
        ('단타 타점 3가지로',     '차트에서 진입점을 찾아봤다'),
        ('짝이랑 함께',          '혼자가 아닌 검증 경험을 했다'),
    ]
    y = Inches(1.9)
    for label, val in summary:
        txt(sl, label, Inches(0.5), y, Inches(5.8), Inches(0.55), size=18, color=GRAY)
        txt(sl, val,   Inches(6.3), y, Inches(6.5), Inches(0.55), size=18, bold=True, color=LIME)
        y += Inches(0.75)

    # 과제
    box(sl, Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.0), hr('#0D1A00'))
    txt(sl, '과제  —  오늘 타점 후보를 찾은 종목 1개를 기록해오기 (3주차: 타점 심화)',
        Inches(0.7), Inches(4.2), Inches(11.8), Inches(0.7),
        size=19, bold=True, color=LIME)

    # 3주차 예고
    box(sl, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.5), hr('#111111'))
    txt(sl, '3주차 예고  —  좋은 타점을 잡는 법 (심화)',
        Inches(0.7), Inches(5.42), Inches(11.5), Inches(0.55),
        size=20, bold=True, color=WHITE)
    txt(sl, '이동평균선, 볼린저밴드, 일목균형표 — 차트 지표로 타점의 신뢰도 높이기',
        Inches(0.7), Inches(5.92), Inches(11.5), Inches(0.55),
        size=16, color=GRAY)

    txt(sl, '피스메이커 자가점검  |  차이보존 · 공감≠동감 · 최소한의힘',
        Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
        size=12, color=hr('#444444'))

# ─────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────
prs = new_prs()

slide_title(prs)
slide_agenda(prs)
slide_opening(prs)
slide_edu1_intro(prs)
slide_scoring_table(prs)
slide_grade(prs)
slide_example(prs)
slide_pair1(prs)
slide_danta_intro(prs)
slide_danta_chart(prs)
slide_danta_rules(prs)
slide_pair2(prs)
slide_closing(prs)

out = r'C:\haessje\라임아카데미\돈복사_2주차.pptx'
prs.save(out)
print(f'완료: {out}  ({len(prs.slides)}슬라이드)')
